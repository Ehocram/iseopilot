"""
Editing IN PLACE dei documenti ALLEGATI (Incremento 10).

Differenza sostanziale rispetto a docgen: qui NON si genera un documento nuovo
dal template, si apre il FILE CARICATO dall'utente e si modificano le sole
parti richieste, restituendo una COPIA con tutto il resto invariato (stili,
intestazioni, loghi, numerazione, tabelle, formule).

Pattern architetturale di sempre: il motore è DETERMINISTICO, l'AI non tocca
il file. Il modello produce solo un PIANO di modifiche ancorate a testo esatto
("trova" → "sostituisci"); il codice Python applica il piano al documento vero
e, quando un ancoraggio non esiste, lo DICHIARA invece di inventare.

Tracciatura per formato (dichiarata all'utente, mai implicita):
  • .docx → REVISIONI TRACCIATE native (w:ins / w:del): chi riceve il file
    apre Word e vede Revisioni, con accetta/rifiuta per singola modifica.
  • .pptx → il formato NON prevede revisioni: modifiche applicate e registrate
    nelle NOTE del relatore della slide interessata.
  • .xlsx → valori aggiornati e vecchio valore conservato in un COMMENTO di
    cella. Rifiuto esplicito se il foglio contiene grafici o immagini, che la
    libreria non è in grado di preservare (non-distruttivo per default).
"""
import copy
import datetime
import json
import os
import re
import tempfile
import zipfile
from pathlib import Path

from . import orchestrator

OUT_DIR = Path(tempfile.gettempdir()) / "iseopilot_docs"
OUT_DIR.mkdir(parents=True, exist_ok=True)

EDITABLE_EXT = {".docx", ".pptx", ".xlsx"}
MAX_MODIFICHE = 40           # tetto per piano: oltre, si dichiara il taglio
AUTORE_REVISIONE = "ISEOPilot"
TESTO_MAX_CHARS = 60000      # testo del documento passato al modello


# ── 1) Rilevamento della richiesta di modifica ──────────────
# Verbi di REVISIONE (non di creazione): il documento esiste già ed è allegato.
_EDIT_VERB_RE = re.compile(
    r"\b(modific\w+|cambia\w*|aggiorn\w+|corregg\w+|sistem\w+|riscriv\w+|"
    r"rived\w+|revision\w+|sostitu\w+|rimpiazz\w+|togli\w*|rimuov\w+|elimin\w+|"
    r"cancell\w+|aggiung\w+|integra\w*|inserisc\w+|accorcia\w*|allung\w+|"
    r"update|modify|change|correct|revise|replace|remove|delete|add|insert|"
    r"shorten|expand|fix)\b", re.IGNORECASE)
# Riferimento al documento allegato
_DOC_REF_RE = re.compile(
    r"\b(document\w*|file|allegat\w+|word|excel|powerpoint|presentazion\w*|"
    r"slide|foglio|cartella\s+di\s+lavoro|tabella|report|contratto|offert\w+|"
    r"proposta|specific\w*|manuale|guida|testo|attachment|document|"
    r"spreadsheet|workbook|deck)\b", re.IGNORECASE)
# Forme clitiche riferite al documento ("modificalo", "correggimelo")
_EDIT_CLITIC_RE = re.compile(
    r"\b(modifical[oa]|correggil[oa]|correggimel[oa]|aggiornal[oa]|"
    r"sistemal[oa]|riscrivil[oa]|rivedil[oa]|revisional[oa]|"
    r"(edit|fix|update|revise)\s+it)\b", re.IGNORECASE)
# Verbi di CREAZIONE: se presenti vince la generazione (docgen), non l'editing
_CREATE_VERB_RE = re.compile(
    r"\b(crea\w*|genera\w*|prepara\w*|produci|redigi|stendi|fammi|farmi|"
    r"create|generate|prepare|draft|make\s+me|build)\b", re.IGNORECASE)
# Domande sul contenuto: non sono ordini di modifica
_EXPLAIN_RE = re.compile(
    r"\b(cosa|cos'è|che\s+cosa|come|perch[éè]|quando|quali|quanto|spiegam\w+|"
    r"riassum\w+|sintetizz\w+|dimmi|elenca|analizz\w+|confront\w+|"
    r"what|how|why|when|which|explain|summar\w+|tell\s+me|list)\b",
    re.IGNORECASE)
# Restituzione esplicita del file ("restituiscimelo", "rimandamelo modificato")
_RETURN_RE = re.compile(
    r"\b(restitu\w+|rimanda\w*|ridamm\w+|ridar\w+|riman\w+|scaricar\w+|"
    r"download|give\s+it\s+back|send\s+it\s+back|return)\b", re.IGNORECASE)


def editable_attachments(attachments: list | None) -> list:
    """Allegati di formato modificabile presenti nel messaggio."""
    out = []
    for a in (attachments or [])[:20]:
        if not a.get("ok") or not a.get("id"):
            continue
        ext = Path(str(a.get("name") or "")).suffix.lower()
        if ext in EDITABLE_EXT:
            out.append(a)
    return out


def detect_edit_request(text: str, attachments: list | None) -> dict | None:
    """Ritorna l'ALLEGATO da modificare in place, oppure None.

    Deterministico e conservativo: senza un allegato modificabile non scatta
    mai; se la richiesta è di CREARE un documento nuovo vince docgen; se è una
    domanda sul contenuto non scatta. In presenza di più allegati modificabili
    si sceglie quello nominato nel messaggio, altrimenti l'ultimo caricato.
    """
    cand = editable_attachments(attachments)
    if not cand:
        return None
    t = (text or "").strip()
    if not t or len(t) > 2000:
        return None
    # una richiesta di CREAZIONE resta di competenza di docgen
    if _CREATE_VERB_RE.search(t):
        return None
    # allegato NOMINATO nel messaggio ("modifica il Listino"): il nome del file
    # è esso stesso un riferimento al documento, quanto la parola "documento".
    tl = t.lower()
    nominato = None
    for a in cand:
        nome = str(a.get("name") or "").lower()
        stem = Path(nome).stem
        if nome and (nome in tl or (len(stem) >= 4 and stem in tl)):
            nominato = a
            break
    ordine = bool(_EDIT_CLITIC_RE.search(t)) or bool(
        _EDIT_VERB_RE.search(t)
        and (_DOC_REF_RE.search(t) or _RETURN_RE.search(t) or nominato))
    if not ordine:
        return None
    # domanda sul contenuto senza richiesta di restituzione del file
    if _EXPLAIN_RE.search(t) and not (_RETURN_RE.search(t) or _EDIT_CLITIC_RE.search(t)):
        return None
    return nominato or cand[-1]


# ── 2) Piano di modifiche (AI: propone, non esegue) ─────────
_SCHEMA_COMUNE = """Rispondi SOLO con JSON valido, senza testo prima o dopo, in questo schema:
{"modifiche": [{"trova": "<testo ESATTO copiato dal documento>",
                "sostituisci": "<nuovo testo, oppure stringa vuota per cancellare>",
                "motivo": "<perché, in una riga>"}],
 "note": "<eventuali osservazioni, o stringa vuota>"}

REGOLE VINCOLANTI:
- "trova" deve essere copiato ALLA LETTERA dal documento fornito, comprese
  maiuscole e punteggiatura: è l'ancoraggio con cui il file verrà modificato.
  Se non è letterale, la modifica NON verrà applicata.
- "trova" deve stare dentro UN SOLO paragrafo (o una sola cella) e deve essere
  abbastanza lungo da essere univoco nel documento.
- Applica SOLO ciò che l'utente ha chiesto: non riscrivere parti non richieste,
  non correggere ciò che non è oggetto della richiesta.
- Se la richiesta non è applicabile al documento, ritorna "modifiche": [] e
  spiega il motivo in "note"."""

_SCHEMA_XLSX = """Rispondi SOLO con JSON valido, senza testo prima o dopo, in questo schema:
{"modifiche": [{"foglio": "<nome foglio, opzionale>",
                "cella": "<riferimento tipo B7, opzionale>",
                "trova": "<valore ESATTO della cella, se non usi 'cella'>",
                "sostituisci": "<nuovo valore>",
                "motivo": "<perché, in una riga>"}],
 "note": "<eventuali osservazioni, o stringa vuota>"}

REGOLE VINCOLANTI:
- Indica "cella" (con "foglio") quando la posizione è nota dal documento;
  altrimenti usa "trova" con il valore ESATTO della cella da cambiare.
- Non toccare celle che contengono FORMULE se non esplicitamente richiesto.
- Applica SOLO ciò che l'utente ha chiesto.
- Se la richiesta non è applicabile, ritorna "modifiche": [] e spiega in "note"."""


def build_plan(fmt: str, richiesta: str, testo_doc: str, settings: dict,
               note_utente: str = "") -> dict:
    """Chiede al modello il PIANO di modifiche. Fail-loud sul JSON invalido."""
    schema = _SCHEMA_XLSX if fmt == "xlsx" else _SCHEMA_COMUNE
    lang = settings.get("reply_lang") or "Italiano"
    system = (
        "Sei un revisore di documenti aziendali ISEO. Non riscrivi il documento: "
        "produci un piano di modifiche MIRATE che verrà applicato da un programma. "
        f"Le modifiche vanno scritte nella lingua del documento; le note in {lang}, "
        "salvo diversa lingua chiesta esplicitamente dall'utente. "
        f"Massimo {MAX_MODIFICHE} modifiche.\n\n" + schema)
    note = ("\n\n" + note_utente.strip()) if note_utente and note_utente.strip() else ""
    user = (f"RICHIESTA DI MODIFICA:\n{richiesta}{note}\n\n"
            f"CONTENUTO ATTUALE DEL DOCUMENTO:\n{testo_doc[:TESTO_MAX_CHARS]}")
    raw = orchestrator.complete(system, user, settings, max_tokens=4000)
    plan = _parse_json(raw)
    mods = plan.get("modifiche") or []
    if not isinstance(mods, list):
        raise ValueError("Piano di modifiche non valido: 'modifiche' non è una lista.")
    plan["troncato"] = len(mods) > MAX_MODIFICHE
    plan["modifiche"] = mods[:MAX_MODIFICHE]
    return plan


def _parse_json(raw: str) -> dict:
    t = (raw or "").strip()
    if t.startswith("```"):
        t = re.sub(r"^```[a-zA-Z]*\s*", "", t)
        t = re.sub(r"\s*```$", "", t)
    i, j = t.find("{"), t.rfind("}")
    if i < 0 or j <= i:
        raise ValueError("Il modello non ha prodotto un piano di modifiche leggibile.")
    return json.loads(t[i:j + 1])


# ── 3) Applicazione DETERMINISTICA al documento ─────────────
def _now_iso() -> str:
    return datetime.datetime.now(datetime.timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


def _norm(s: str) -> str:
    """Normalizzazione tollerante per l'ancoraggio: spazi unificati, virgolette
    tipografiche ricondotte a quelle dritte. NON altera il documento: serve solo
    a ritrovare il punto quando il modello ricopia con spaziature diverse."""
    s = (s or "").replace("\u00a0", " ").replace("\u2019", "'").replace("\u2018", "'")
    s = s.replace("\u201c", '"').replace("\u201d", '"')
    s = s.replace("\u2013", "-").replace("\u2014", "-")
    return re.sub(r"\s+", " ", s).strip()


# ---- DOCX: revisioni tracciate native ----
def _mk_run(template_r, text: str, deleted: bool = False):
    """Run clonato dal run originale (ne eredita la formattazione) con nuovo
    testo. Nei blocchi <w:del> il testo va in <w:delText>, non in <w:t>."""
    from docx.oxml.ns import qn
    r = copy.deepcopy(template_r)
    for tag in ("w:t", "w:delText"):
        for el in r.findall(qn(tag)):
            r.remove(el)
    el = r.makeelement(qn("w:delText") if deleted else qn("w:t"),
                       {qn("xml:space"): "preserve"})
    el.text = text
    r.append(el)
    return r


def _docx_paragrafi(doc):
    """Paragrafi del corpo e delle tabelle (anche annidate)."""
    for p in doc.paragraphs:
        yield p
    def _tab(tables):
        for t in tables:
            for row in t.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        yield p
                    yield from _tab(cell.tables)
    yield from _tab(doc.tables)


def _docx_applica_in_paragrafo(p, trova: str, sostituisci: str, rid: int) -> bool:
    """Applica UNA modifica dentro un paragrafo come revisione tracciata.
    Ritorna True se l'ancoraggio è stato trovato e la modifica inserita."""
    from docx.oxml.ns import qn
    runs = list(p.runs)
    if not runs:
        return False
    full = "".join(r.text or "" for r in runs)
    idx = full.find(trova)
    if idx < 0:                       # secondo tentativo: confronto normalizzato
        nf, nt = _norm(full), _norm(trova)
        if not nt or nt not in nf:
            return False
        # mappa la posizione normalizzata sull'originale, carattere per carattere
        mapping, buf, prev_space = [], [], False
        for i, ch in enumerate(full):
            c = _norm(ch) if ch.strip() else " "
            if c == " ":
                if prev_space or not buf:
                    continue
                prev_space = True
                buf.append(" "); mapping.append(i)
            else:
                prev_space = False
                buf.append(c); mapping.append(i)
        norm_join = "".join(buf).strip()
        off = len("".join(buf)) - len("".join(buf).lstrip())
        pos = norm_join.find(nt)
        if pos < 0:
            return False
        idx = mapping[pos + off]
        fine = mapping[min(pos + off + len(nt) - 1, len(mapping) - 1)] + 1
    else:
        fine = idx + len(trova)

    # run che intersecano l'intervallo [idx, fine)
    over, pos = [], 0
    for r in runs:
        s = pos
        e = pos + len(r.text or "")
        pos = e
        if e > idx and s < fine:
            over.append((r, s, e))
    if not over:
        return False

    primo_r = over[0][0]._r
    pre = full[over[0][1]:idx]
    post = full[fine:over[-1][2]]
    attrs = {qn("w:id"): str(rid), qn("w:author"): AUTORE_REVISIONE,
             qn("w:date"): _now_iso()}

    nuovi = []
    if pre:
        nuovi.append(_mk_run(primo_r, pre))
    del_el = primo_r.makeelement(qn("w:del"), attrs)
    for r, s, e in over:
        seg = full[max(s, idx):min(e, fine)]
        if seg:
            del_el.append(_mk_run(r._r, seg, deleted=True))
    if len(del_el):
        nuovi.append(del_el)
    if sostituisci:
        ins_el = primo_r.makeelement(qn("w:ins"), dict(attrs, **{qn("w:id"): str(rid + 1)}))
        ins_el.append(_mk_run(primo_r, sostituisci))
        nuovi.append(ins_el)
    if post:
        nuovi.append(_mk_run(primo_r, post))

    for el in nuovi:
        primo_r.addprevious(el)
    for r, _s, _e in over:
        r._r.getparent().remove(r._r)
    return True


def apply_docx(src: str, plan: dict) -> tuple:
    """Applica il piano a un .docx come REVISIONI TRACCIATE.
    Ritorna (path_output, applicate, non_applicate, avvisi)."""
    from docx import Document
    doc = Document(src)
    applicate, non_applicate = [], []
    rid = 1000
    for m in plan.get("modifiche", []):
        trova = str(m.get("trova") or "")
        sost = str(m.get("sostituisci") or "")
        if not trova:
            non_applicate.append({**m, "perche": "ancoraggio 'trova' assente nel piano"})
            continue
        fatto = False
        for p in _docx_paragrafi(doc):
            if _docx_applica_in_paragrafo(p, trova, sost, rid):
                rid += 2
                fatto = True
                break
        (applicate if fatto else non_applicate).append(
            m if fatto else {**m, "perche": "testo non trovato nel documento"})
    out = str(OUT_DIR / f"edit_{os.urandom(6).hex()}.docx")
    doc.save(out)
    return out, applicate, non_applicate, []


# ---- PPTX: applicazione diretta + registro nelle note ----
def _pptx_applica_in_paragrafo(par, trova: str, sostituisci: str) -> bool:
    runs = list(par.runs)
    if not runs:
        return False
    full = "".join(r.text or "" for r in runs)
    idx = full.find(trova)
    if idx < 0:
        nf, nt = _norm(full), _norm(trova)
        if not nt or nt not in nf:
            return False
        idx = full.lower().find(trova[:20].lower())
        if idx < 0:
            return False
        fine = min(idx + len(trova), len(full))
    else:
        fine = idx + len(trova)
    pos, primo = 0, True
    for r in runs:
        s = pos
        e = pos + len(r.text or "")
        pos = e
        if e <= idx or s >= fine:
            continue
        pre = (r.text or "")[:max(0, idx - s)]
        post = (r.text or "")[max(0, fine - s):]
        r.text = pre + (sostituisci if primo else "") + post
        primo = False
    return True


def apply_pptx(src: str, plan: dict) -> tuple:
    """Applica il piano a un .pptx. Il formato NON prevede revisioni tracciate:
    le modifiche sono applicate e REGISTRATE nelle note del relatore."""
    from pptx import Presentation
    prs = Presentation(src)
    applicate, non_applicate = [], []
    registro = {}
    for m in plan.get("modifiche", []):
        trova = str(m.get("trova") or "")
        sost = str(m.get("sostituisci") or "")
        if not trova:
            non_applicate.append({**m, "perche": "ancoraggio 'trova' assente nel piano"})
            continue
        fatto = False
        for n, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for par in shape.text_frame.paragraphs:
                    if _pptx_applica_in_paragrafo(par, trova, sost):
                        registro.setdefault(n, []).append((trova, sost))
                        fatto = True
                        break
                if fatto:
                    break
            if fatto:
                break
        (applicate if fatto else non_applicate).append(
            m if fatto else {**m, "perche": "testo non trovato nelle slide"})
    # registro non distruttivo nelle note del relatore
    for n, voci in registro.items():
        try:
            slide = prs.slides[n - 1]
            tf = slide.notes_slide.notes_text_frame
            righe = [f"[{AUTORE_REVISIONE} — modifiche applicate il "
                     f"{datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}]"]
            for trova, sost in voci:
                righe.append(f"• «{trova[:120]}» → «{sost[:120]}»" if sost
                             else f"• rimosso: «{trova[:120]}»")
            tf.text = (tf.text + "\n" if tf.text else "") + "\n".join(righe)
        except Exception:
            pass          # le note sono un di più: non compromettono il file
    out = str(OUT_DIR / f"edit_{os.urandom(6).hex()}.pptx")
    prs.save(out)
    return out, applicate, non_applicate, [
        "PowerPoint non prevede revisioni tracciate: le modifiche sono applicate "
        "e riepilogate nelle note del relatore delle slide interessate."]


# ---- XLSX: valori aggiornati + vecchio valore in commento ----
def xlsx_ha_grafici_o_immagini(src: str) -> bool:
    """True se la cartella contiene grafici/immagini/pivot: la libreria di
    scrittura NON è in grado di preservarli e il salvataggio li perderebbe."""
    try:
        with zipfile.ZipFile(src) as z:
            for n in z.namelist():
                nl = n.lower()
                if (nl.startswith("xl/charts/") or nl.startswith("xl/media/")
                        or nl.startswith("xl/drawings/drawing")
                        or nl.startswith("xl/pivotcache/")):
                    return True
    except Exception:
        return False
    return False


def apply_xlsx(src: str, plan: dict) -> tuple:
    """Applica il piano a un .xlsx: valore aggiornato e VECCHIO valore
    conservato in un commento di cella."""
    import openpyxl
    from openpyxl.comments import Comment
    wb = openpyxl.load_workbook(src)
    applicate, non_applicate = [], []
    for m in plan.get("modifiche", []):
        sost = m.get("sostituisci")
        sost = "" if sost is None else str(sost)
        foglio = str(m.get("foglio") or "").strip()
        cella = str(m.get("cella") or "").strip().upper()
        trova = str(m.get("trova") or "")
        target = None
        if cella:
            ws = wb[foglio] if foglio and foglio in wb.sheetnames else wb[wb.sheetnames[0]]
            try:
                target = (ws, ws[cella])
            except Exception:
                target = None
        if target is None and trova:
            fogli = [wb[foglio]] if foglio and foglio in wb.sheetnames else list(wb.worksheets)
            for ws in fogli:
                for row in ws.iter_rows():
                    for c in row:
                        if c.value is not None and _norm(str(c.value)) == _norm(trova):
                            target = (ws, c)
                            break
                    if target:
                        break
                if target:
                    break
        if target is None:
            non_applicate.append({**m, "perche": "cella o valore non trovati nel foglio"})
            continue
        ws, c = target
        vecchio = "" if c.value is None else str(c.value)
        if isinstance(c.value, str) and c.value.startswith("="):
            # Non-distruttivo: una formula sostituita da una costante è un danno
            # difficile da accorgersi. Si dichiara e si lascia intatta, sempre.
            non_applicate.append({
                **m, "cella": f"{ws.title}!{c.coordinate}",
                "perche": f"la cella contiene una formula ({vecchio[:60]}): "
                          "non modificata, intervieni a mano se necessario"})
            continue
        try:
            c.value = _coerce(sost)
            c.comment = Comment(f"{AUTORE_REVISIONE}: valore precedente «{vecchio[:200]}»",
                                AUTORE_REVISIONE)
            applicate.append({**m, "cella": f"{ws.title}!{c.coordinate}", "precedente": vecchio})
        except Exception as e:
            non_applicate.append({**m, "perche": f"cella non scrivibile: {str(e)[:80]}"})
    out = str(OUT_DIR / f"edit_{os.urandom(6).hex()}.xlsx")
    wb.save(out)
    return out, applicate, non_applicate, [
        "Excel non prevede revisioni tracciate: i valori sono aggiornati e il "
        "valore precedente è conservato in un commento di cella."]


def _coerce(v: str):
    """Numeri e date restano tali; il resto è testo."""
    t = (v or "").strip()
    if not t:
        return None
    try:
        if re.fullmatch(r"-?\d+", t):
            return int(t)
        if re.fullmatch(r"-?\d+[.,]\d+", t):
            return float(t.replace(",", "."))
    except Exception:
        pass
    return v


_APPLICATORI = {"docx": apply_docx, "pptx": apply_pptx, "xlsx": apply_xlsx}


def apply_plan(fmt: str, src: str, plan: dict) -> tuple:
    if fmt not in _APPLICATORI:
        raise ValueError(f"Formato non modificabile: {fmt}")
    if fmt == "xlsx" and xlsx_ha_grafici_o_immagini(src):
        raise ValueError(
            "Il foglio contiene grafici, immagini o tabelle pivot che la modifica "
            "automatica non è in grado di preservare: il file NON è stato toccato. "
            "Posso elencarti le modifiche da applicare a mano, oppure procedere su "
            "una versione del foglio senza grafici.")
    return _APPLICATORI[fmt](src, plan)


def riepilogo(nome_file: str, fmt: str, applicate: list, non_applicate: list,
              avvisi: list, troncato: bool = False, note: str = "") -> str:
    """Riepilogo in chat: cosa è stato applicato e cosa NO, sempre entrambi."""
    r = []
    if fmt == "docx":
        r.append(f"Ho modificato **{nome_file}** con **revisioni tracciate**: "
                 "aprendolo in Word trovi le modifiche in Revisioni, da accettare "
                 "o rifiutare una per una.")
    else:
        r.append(f"Ho modificato **{nome_file}**.")
    if applicate:
        r.append(f"\n**Modifiche applicate ({len(applicate)}):**")
        for m in applicate[:20]:
            dove = f" `{m['cella']}`" if m.get("cella") else ""
            motivo = f" — {m['motivo']}" if m.get("motivo") else ""
            sost = str(m.get("sostituisci") or "")
            trova = str(m.get("trova") or "")
            if sost:
                r.append(f"- {dove.strip() or '«' + trova[:60] + '»'} → «{sost[:80]}»{motivo}")
            else:
                r.append(f"- rimosso: «{trova[:60]}»{motivo}")
        if len(applicate) > 20:
            r.append(f"- …e altre {len(applicate) - 20}.")
    else:
        r.append("\n**Nessuna modifica applicata.**")
    if non_applicate:
        r.append(f"\n**NON applicate ({len(non_applicate)}) — da verificare a mano:**")
        for m in non_applicate[:10]:
            r.append(f"- «{str(m.get('trova') or '')[:60]}» — {m.get('perche', 'motivo non noto')}")
        if len(non_applicate) > 10:
            r.append(f"- …e altre {len(non_applicate) - 10}.")
    if troncato:
        r.append(f"\n*Piano troncato al tetto di {MAX_MODIFICHE} modifiche: "
                 "rilancia sulle parti rimanenti.*")
    for a in avvisi:
        r.append(f"\n*{a}*")
    if note:
        r.append(f"\n*Nota: {note[:400]}*")
    r.append("\n*L'originale che hai caricato non è stato toccato: questa è una copia.*")
    return "\n".join(r)
