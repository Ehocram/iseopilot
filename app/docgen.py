"""Generazione di file scaricabili (Word, PowerPoint, Excel, PDF).

Flusso: l'utente chiede in chat "creami un word/excel/ppt/pdf …". Qui si rileva
il formato, si fa generare a Claude il CONTENUTO in JSON strutturato, e si
costruisce il file con le librerie sul server. Word e PowerPoint usano i TEMPLATE
aziendali (templateIseo.docx, Presentation_template_1.pptx); Excel è generato da
zero. Per il PDF si converte il Word (LibreOffice se presente) con fallback a un
PDF in stile ISEO.

I file finiti vanno in una cartella temporanea e vengono serviti via token
(vedi connectors.register_report / main download endpoint).
"""
from __future__ import annotations

import json
import os
import re
import tempfile
from pathlib import Path

from . import orchestrator

TEMPLATES_DIR = Path(__file__).resolve().parent / "doc_templates"
DOCX_TEMPLATE = TEMPLATES_DIR / "templateIseo.docx"
PPTX_TEMPLATE = TEMPLATES_DIR / "Presentation_template_1.pptx"
OUT_DIR = Path(tempfile.gettempdir()) / "iseopilot_docs"
OUT_DIR.mkdir(parents=True, exist_ok=True)

# ── Template PERSONALI per-utente (bypass del template di default) ──
# L'utente può caricare dal composer un proprio .docx/.pptx: finché è
# presente, la generazione usa QUELLO al posto del template ISEO. Storage
# isolato per identità (stesso hashing dei token), niente macro, audit a
# monte (main.py). Un solo slot per formato per utente: stato visibile in
# chat come chip, quindi nessuna sorpresa tra una conversazione e l'altra.
USER_TPL_DIR = Path(os.environ.get("APP_DATA_DIR", "/data")) / "user_templates"
TPL_MAX_BYTES = 15 * 1024 * 1024
_TPL_EXT = {"docx": ".docx", "pptx": ".pptx"}


def _tpl_user_dir(user: str) -> Path:
    import hashlib
    d = USER_TPL_DIR / hashlib.sha256(user.encode()).hexdigest()[:16]
    d.mkdir(parents=True, exist_ok=True)
    return d


def validate_office_template(raw: bytes, fmt: str, filename: str = "") -> str:
    """Valida un template caricato. Ritorna "" se valido, altrimenti il MOTIVO
    del rifiuto (parlante, per l'utente). Regole: solo .docx/.pptx, archivio
    Office leggibile con le part attese, NIENTE macro (macroEnabled/vbaProject)."""
    import io
    import zipfile
    name_l = (filename or "").lower()
    if name_l.endswith((".docm", ".pptm", ".dotm", ".potm", ".xlsm")):
        return "I template con macro non sono ammessi (.docm/.pptm): usa un .docx o .pptx senza macro."
    if fmt not in _TPL_EXT:
        return "Formato non supportato: sono ammessi solo template .docx (Word) e .pptx (PowerPoint)."
    if len(raw) > TPL_MAX_BYTES:
        return f"File troppo grande ({len(raw)//(1024*1024)} MB): il limite è {TPL_MAX_BYTES//(1024*1024)} MB."
    bio = io.BytesIO(raw)
    if not zipfile.is_zipfile(bio):
        return "Il file non è un documento Office valido (archivio non leggibile)."
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as z:
            names = z.namelist()
            if "[Content_Types].xml" not in names:
                return "Il file non è un documento Office valido ([Content_Types].xml assente)."
            if any("vbaproject" in n.lower() for n in names):
                return "Il template contiene macro (vbaProject): non ammesso."
            ct = z.read("[Content_Types].xml").decode("utf-8", "replace").lower()
            if "macroenabled" in ct:
                return "Il template è macro-enabled: non ammesso."
            need = "word/document.xml" if fmt == "docx" else "ppt/presentation.xml"
            if need not in names:
                atteso = "Word (.docx)" if fmt == "docx" else "PowerPoint (.pptx)"
                return f"Il file non è un template {atteso} valido ({need} assente)."
    except Exception as e:
        return f"Template non leggibile: {e}"
    return ""


# I MODELLI Word/PowerPoint (.dotx/.potx, senza macro) sono lo stesso package
# OPC con un content-type diverso: python-docx/python-pptx li rifiutano. La
# normalizzazione riscrive SOLO l'override del part principale in
# [Content_Types].xml: il resto (stili, header, footer, tema, numbering)
# resta byte-identico.
_CT_TEMPLATE_TO_DOC = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.template.main+xml":
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
    "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml":
        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
}


def normalize_office_template(raw: bytes) -> bytes:
    """Se il package è un MODELLO (.dotx/.potx), lo converte in documento
    (.docx/.pptx) riscrivendo il content-type del part principale. Se è già
    un documento, ritorna i byte invariati. In caso di errore, invariati."""
    import io
    import zipfile
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as z:
            ct = z.read("[Content_Types].xml").decode("utf-8", "replace")
            hit = None
            for tpl_ct, doc_ct in _CT_TEMPLATE_TO_DOC.items():
                if tpl_ct in ct:
                    hit = (tpl_ct, doc_ct)
                    break
            if not hit:
                return raw
            out = io.BytesIO()
            with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zo:
                for item in z.infolist():
                    data = z.read(item.filename)
                    if item.filename == "[Content_Types].xml":
                        data = ct.replace(hit[0], hit[1]).encode("utf-8")
                    zo.writestr(item, data)
            return out.getvalue()
    except Exception:
        return raw


def save_user_template(user: str, fmt: str, raw: bytes, orig_name: str) -> str:
    """Salva il template dell'utente per il formato. Ritorna "" o l'errore."""
    err = validate_office_template(raw, fmt, orig_name)
    if err:
        return err
    raw = normalize_office_template(raw)  # .dotx/.potx -> .docx/.pptx
    d = _tpl_user_dir(user)
    (d / f"template{_TPL_EXT[fmt]}").write_bytes(raw)
    meta = {}
    mp = d / "meta.json"
    try:
        if mp.exists():
            meta = json.loads(mp.read_text(encoding="utf-8"))
    except Exception:
        meta = {}
    meta[fmt] = (orig_name or f"template{_TPL_EXT[fmt]}")[:120]
    mp.write_text(json.dumps(meta, ensure_ascii=False), encoding="utf-8")
    return ""


def get_user_template(user: str, fmt: str):
    """(path, nome_originale) del template dell'utente, oppure None."""
    d = _tpl_user_dir(user)
    p = d / f"template{_TPL_EXT.get(fmt, '')}"
    if not p.is_file():
        return None
    nome = f"template{_TPL_EXT[fmt]}"
    try:
        meta = json.loads((d / "meta.json").read_text(encoding="utf-8"))
        nome = meta.get(fmt, nome)
    except Exception:
        pass
    return str(p), nome


def delete_user_template(user: str, fmt: str) -> bool:
    d = _tpl_user_dir(user)
    p = d / f"template{_TPL_EXT.get(fmt, '')}"
    if p.is_file():
        try:
            p.unlink()
        except Exception:
            return False
    try:
        mp = d / "meta.json"
        if mp.exists():
            meta = json.loads(mp.read_text(encoding="utf-8"))
            meta.pop(fmt, None)
            mp.write_text(json.dumps(meta, ensure_ascii=False), encoding="utf-8")
    except Exception:
        pass
    return True


def user_templates_status(user: str) -> dict:
    out = {}
    for fmt in ("docx", "pptx"):
        t = get_user_template(user, fmt)
        out[fmt] = {"name": t[1]} if t else None
    return out

ISEO_RED = "EC1D2B"

# ── 1) Rilevamento intento + formato ───────────────────────
_FMT_WORD = ("word", "documento word", ".docx", "in word", "un doc ")
_FMT_XLS = ("excel", "foglio di calcolo", "foglio elettronico", "spreadsheet", ".xlsx", "in excel")
_FMT_PPT = ("powerpoint", "power point", "presentazione", "slide", "slides", "ppt", ".pptx", "deck")
_FMT_PDF = ("pdf", ".pdf", "in pdf")
# Verbi di generazione: forme FORTI (creazione esplicita, con le coniugazioni
# reali: "mi generi", "mi crei", "fammi") e forme DEBOLI (cortesia/desiderio).
# Le deboli valgono solo se la frase non è una richiesta di spiegazione
# ("puoi spiegarmi come funziona Excel?" resta una chat, non genera file).
_STRONG_VERB_RE = re.compile(
    r"\b(crea(re|mi|temi)?|crei|genera(re|mi)?|generi|fammi|fai|fatemi|"
    r"produci|produrre|produrmi|prepara(re|mi)?|prepari|esporta(re|mi)?|esporti|"
    r"scarica(re|mi)?|scarichi|redigi|redigere|stila(re)?|stilami|"
    r"realizza(re|mi)?|realizzi|costruisci(mi)?|costruire|buttami|"
    r"create|generate|make|build|produce|export|download|prepare|draft)\b",
    re.IGNORECASE)
_WEAK_VERB_RE = re.compile(
    r"\b(vorrei|voglio|desidero|mi\s+serve|servirebbe|avrei\s+bisogno|"
    r"puoi|potresti|riesci|sapresti|"
    r"can\s+you|could\s+you|i\s+need|i\s+want|i'd\s+like)\b",
    re.IGNORECASE)
_EXPLAIN_RE = re.compile(
    r"\b(spiega|spiegami|spieghi|cos'?\s?è|cosa\s+è|come\s+funziona|"
    r"come\s+si\s+usa|aiutami\s+a\s+capire|che\s+cos|differenza\s+tra|"
    r"explain|what\s+is|how\s+does|how\s+to\s+use)\b",
    re.IGNORECASE)


def detect_request(text: str) -> str | None:
    """Ritorna il formato richiesto ('docx'|'xlsx'|'pptx'|'pdf') o None."""
    tl = (text or "").lower()
    # l'ordine conta: pptx e xlsx prima di docx/pdf generici
    if any(k in tl for k in _FMT_PPT):
        fmt = "pptx"
    elif any(k in tl for k in _FMT_XLS):
        fmt = "xlsx"
    elif any(k in tl for k in _FMT_WORD):
        fmt = "docx"
    elif any(k in tl for k in _FMT_PDF):
        fmt = "pdf"
    else:
        return None
    # estensione esplicita: intento inequivocabile
    if any(e in tl for e in (".docx", ".xlsx", ".pptx", ".pdf")):
        return fmt
    if _STRONG_VERB_RE.search(tl):
        return fmt
    if _WEAK_VERB_RE.search(tl) and not _EXPLAIN_RE.search(tl):
        return fmt
    return None


# Conferme brevi ("sì procedi", "ok vai", "genera pure"): non contengono né
# formato né verbo, ma se la conversazione stava preparando un documento
# devono far scattare la generazione ereditando il formato dai turni recenti.
_CONFIRM_RE = re.compile(
    r"^\s*(s[iì]|ok|okay|va\s+bene|perfetto|d'accordo|certo|esatto|conferm\w*|"
    r"procedi(amo)?|vai|prosegui|genera(lo|la)?\s*pure|fallo|falla|crealo|creala|"
    r"yes|sure|go\s+ahead|proceed|do\s+it|please\s+do)\b[\s!,.]*",
    re.IGNORECASE)


def detect_request_with_history(text: str, history: list | None = None) -> str | None:
    """Come detect_request, ma se il messaggio è una CONFERMA breve eredita il
    formato dai turni recenti della conversazione (utente o assistente)."""
    fmt = detect_request(text)
    if fmt:
        return fmt
    tl = (text or "").strip()
    if not history or len(tl) > 60:
        return None
    m = _CONFIRM_RE.match(tl)
    if not m:
        return None
    # la conferma deve essere (quasi) solo conferma: "sì ma prima spiegami…"
    # non è un via libera alla generazione
    rest = tl[m.end():].strip(" !,.?")
    if rest and not _CONFIRM_RE.match(rest) and len(rest) > 12:
        return None
    # cerca il formato più recente citato in conversazione (ultimi 6 turni)
    for m in reversed(list(history)[-6:]):
        f = detect_request(str(m.get("content", "")))
        if not f:
            # anche un semplice riferimento al formato nel turno assistente
            # ("procedo con la generazione del file PowerPoint") conta
            c = str(m.get("content", "")).lower()
            for kws, ff in ((_FMT_PPT, "pptx"), (_FMT_XLS, "xlsx"),
                            (_FMT_WORD, "docx"), (_FMT_PDF, "pdf")):
                if any(k in c for k in kws):
                    f = ff
                    break
        if f:
            return f
    return None


# ── 2) Generazione del contenuto strutturato (Claude → JSON) ─
def _parse_json(raw: str) -> dict:
    raw = (raw or "").strip()
    raw = re.sub(r"^```(?:json)?", "", raw).strip()
    raw = re.sub(r"```$", "", raw).strip()
    # estrai il primo oggetto JSON bilanciato
    start = raw.find("{")
    if start < 0:
        raise ValueError("Nessun JSON nella risposta del modello.")
    depth, end = 0, None
    for i in range(start, len(raw)):
        if raw[i] == "{":
            depth += 1
        elif raw[i] == "}":
            depth -= 1
            if depth == 0:
                end = i + 1
                break
    return json.loads(raw[start:end] if end else raw[start:])


_SCHEMAS = {
    "docx": (
        'Genera SOLO un oggetto JSON valido, senza testo prima o dopo, con questo schema:\n'
        '{"title": "Titolo del documento",\n'
        ' "subtitle": "sottotitolo opzionale",\n'
        ' "sections": [{"heading": "Titolo sezione", "paragraphs": ["par1","par2"], "bullets": ["punto1","punto2"]}]}\n'
        'Ogni sezione può avere "paragraphs" e/o "bullets". Scrivi contenuto reale, completo e professionale.'
    ),
    "pdf": None,  # usa lo schema docx
    "pptx": (
        'Genera SOLO un oggetto JSON valido, senza testo prima o dopo, con questo schema:\n'
        '{"title": "Titolo presentazione", "subtitle": "sottotitolo",\n'
        ' "slides": [{"title": "Titolo slide", "bullets": ["punto1","punto2","punto3"]}]}\n'
        'Da 4 a 12 slide. Ogni slide 3-6 bullet concisi (max ~12 parole). Contenuto reale e professionale.'
    ),
    "xlsx": (
        'Genera SOLO un oggetto JSON valido, senza testo prima o dopo, con questo schema:\n'
        '{"title": "Titolo", "sheets": [{"name": "Foglio1",\n'
        '  "columns": ["Col A","Col B","Importo"],\n'
        '  "rows": [["v1","v2",100],["v3","v4",200]],\n'
        '  "total_columns": [2]}]}\n'
        '"total_columns" (opzionale) = indici (0-based) delle colonne numeriche da totalizzare con una formula SOMMA. '
        'Usa numeri reali (non stringhe) nelle celle numeriche. Contenuto reale e completo.'
    ),
}


def build_spec(fmt: str, user_request: str, context: str, settings: dict,
               history_text: str = "", note_utente: str = "") -> dict:
    schema = _SCHEMAS["docx"] if fmt in ("docx", "pdf") else _SCHEMAS[fmt]
    lang = settings.get("reply_lang", "Italiano")
    system = (
        f"Sei un generatore di documenti aziendali ISEO. Scrivi i contenuti nella "
        f"lingua: {lang}, SALVO che la RICHIESTA DELL'UTENTE o la conversazione "
        "chiedano esplicitamente un'altra lingua: in quel caso PREVALE la lingua "
        "richiesta. "
        "Produci ESCLUSIVAMENTE JSON valido secondo lo schema indicato, senza commenti, "
        "senza markdown, senza testo aggiuntivo. Il contenuto deve essere reale, utile e "
        "PERTINENTE alla richiesta specifica (mai generico), basato sul materiale e sulla "
        "conversazione forniti quando presenti."
    )
    ctx = ("\n\n=== MATERIALE DI RIFERIMENTO ===\n" + context.strip()) if context and context.strip() else ""
    hist = ("\n\n=== CONVERSAZIONE PRECEDENTE (per capire cosa vuole davvero l'utente) ===\n"
            + history_text.strip()) if history_text and history_text.strip() else ""
    note = ("\n\n" + note_utente.strip()) if note_utente and note_utente.strip() else ""
    user = f"{schema}\n\nRICHIESTA DELL'UTENTE:\n{user_request}{note}{hist}{ctx}"
    raw = orchestrator.complete(system, user, settings, max_tokens=4000)
    return _parse_json(raw)


# ── 3) Costruttori per formato ─────────────────────────────
def _safe_name(title: str, ext: str) -> str:
    base = re.sub(r"[^0-9A-Za-zÀ-ÿ _-]+", "", (title or "documento")).strip()[:48] or "documento"
    return f"{base}.{ext}"


def _find_bullet_num_id(doc) -> int | None:
    """Primo numId del template il cui abstractNum ha numFmt='bullet' al
    livello 0: serve per elenchi puntati REALI sui template personali (dove
    il numId=1 del template ISEO non esiste). None se non determinabile."""
    try:
        from docx.oxml.ns import qn
        npart = getattr(doc.part, "numbering_part", None)
        if npart is None:
            return None
        root = npart.element
        bullet_abstracts = set()
        for an in root.findall(qn("w:abstractNum")):
            lvl0 = an.find(qn("w:lvl"))
            fmt = lvl0.find(qn("w:numFmt")) if lvl0 is not None else None
            if fmt is not None and fmt.get(qn("w:val")) == "bullet":
                bullet_abstracts.add(an.get(qn("w:abstractNumId")))
        candidates = []
        for num in root.findall(qn("w:num")):
            aid = num.find(qn("w:abstractNumId"))
            if aid is not None and aid.get(qn("w:val")) in bullet_abstracts:
                candidates.append(int(num.get(qn("w:numId"))))
        return min(candidates) if candidates else None
    except Exception:
        return None


def gen_docx(spec: dict, template_path: str | None = None) -> tuple[str, str]:
    from docx import Document
    from docx.shared import Pt
    from docx.oxml.ns import qn
    if template_path:
        # template PERSONALE dell'utente: se non è apribile l'errore arriva in
        # chat — MAI ripiegare in silenzio sul template di default.
        try:
            doc = Document(str(template_path))
        except Exception as e:
            raise ValueError(f"Il template Word caricato non è utilizzabile: {e}. "
                             "Rimuovilo o sostituiscilo dal composer.") from e
    else:
        doc = Document(str(DOCX_TEMPLATE)) if DOCX_TEMPLATE.is_file() else Document()
    # svuota i paragrafi placeholder del template (mantiene header/footer/logo)
    for p in list(doc.paragraphs):
        p._element.getparent().remove(p._element)

    def has_style(name: str) -> bool:
        try:
            _ = doc.styles[name]
            return True
        except Exception:
            return False

    # NB: add_paragraph(text, style=...) prima INSERISCE il testo e poi assegna
    # lo stile: se lo stile non esiste, l'eccezione arriva DOPO e il paragrafo
    # grezzo resta nel documento — era la causa delle righe DUPLICATE. Qui lo
    # stile viene applicato sullo STESSO paragrafo, con ripiego senza doppioni.
    def add(text: str, style: str | None = None, bold: bool = False,
            before: int = 0, after: int = 0, bullet: bool = False,
            num_id: int = 1):
        p = doc.add_paragraph()
        if style and has_style(style):
            p.style = doc.styles[style]
        run = p.add_run(str(text))
        run.bold = bold
        if before:
            p.paragraph_format.space_before = Pt(before)
        if after:
            p.paragraph_format.space_after = Pt(after)
        if bullet:
            # elenco VERO del template (numId del template, mai "•" letterale)
            numpr = p._p.get_or_add_pPr().makeelement(qn("w:numPr"), {})
            ilvl = numpr.makeelement(qn("w:ilvl"), {qn("w:val"): "0"})
            numid = numpr.makeelement(qn("w:numId"), {qn("w:val"): str(num_id)})
            numpr.append(ilvl); numpr.append(numid)
            p._p.get_or_add_pPr().append(numpr)
        return p

    title_style = "ISEO Titolo" if has_style("ISEO Titolo") else "Title"
    body_style = "Paragrafo base ISEO" if has_style("Paragrafo base ISEO") else "Normal"
    list_style = "Paragrafo elenco" if has_style("Paragrafo elenco") else "List Bullet"
    use_numpr = has_style("Paragrafo elenco") and not template_path
    bullet_num = 1
    subtitle_style, heading_style = body_style, None
    if template_path:
        # TEMPLATE PERSONALE: si parla la lingua del file — Subtitle e
        # Heading 1 se definiti (è lì che vive il branding, es. Solution
        # Centre), List Paragraph/List Bullet per gli elenchi, e il numbering
        # BULLET del template stesso (numId reale, non l'1 del template ISEO).
        if has_style("Subtitle"):
            subtitle_style = "Subtitle"
        if has_style("Heading 1"):
            heading_style = "Heading 1"
        if has_style("List Paragraph"):
            list_style = "List Paragraph"
        elif not has_style("List Bullet"):
            list_style = body_style
        found = _find_bullet_num_id(doc)
        if found is not None:
            use_numpr, bullet_num = True, found

    add(spec.get("title", "Documento"), style=title_style, after=6)
    if spec.get("subtitle"):
        add(spec["subtitle"], style=subtitle_style, after=10)
    for sec in spec.get("sections", []):
        if sec.get("heading"):
            if heading_style:
                # template personale con Heading 1: spaziatura dello stile
                add(sec["heading"], style=heading_style)
            else:
                # titolo di SEZIONE (template ISEO): corpo in grassetto con
                # aria sopra — mai lo stile del titolo di copertina
                add(sec["heading"], style=body_style, bold=True, before=14, after=6)
        for par in sec.get("paragraphs", []):
            add(str(par), style=body_style, after=6)
        for b in sec.get("bullets", []):
            fallback_glyph = (not use_numpr) and bool(template_path)
            add(("• " if fallback_glyph else "") + str(b), style=list_style,
                after=3, bullet=use_numpr, num_id=bullet_num)

    name = _safe_name(spec.get("title", "documento"), "docx")
    path = OUT_DIR / (next(tempfile._get_candidate_names()) + ".docx")
    doc.save(str(path))
    return str(path), name


def gen_pptx(spec: dict, template_path: str | None = None) -> tuple[str, str]:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if template_path:
        try:
            prs = Presentation(str(template_path))
        except Exception as e:
            raise ValueError(f"Il template PowerPoint caricato non è utilizzabile: {e}. "
                             "Rimuovilo o sostituiscilo dal composer.") from e
        if not len(prs.slide_layouts):
            raise ValueError("Il template PowerPoint caricato non contiene layout: "
                             "non utilizzabile per la generazione.")
    else:
        prs = Presentation(str(PPTX_TEMPLATE)) if PPTX_TEMPLATE.is_file() else Presentation()
    # rimuovi le slide di esempio del template (mantiene tema, master, footer):
    # va rimossa sia la voce nell'elenco sia la relazione, altrimenti le part
    # restano orfane e generano nomi duplicati nello zip.
    from pptx.oxml.ns import qn
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        rId = sld_id.get(qn("r:id"))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sld_id_lst.remove(sld_id)

    layout = prs.slide_layouts[0]
    W = prs.slide_width

    def add_textbox(slide, left, top, width, height):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tb.text_frame.word_wrap = True
        return tb

    # copertina
    cover = prs.slides.add_slide(layout)
    t = add_textbox(cover, Inches(0.6), Inches(1.8), W - Inches(1.2), Inches(1.6))
    run = t.text_frame.paragraphs[0].add_run()
    run.text = spec.get("title", "Presentazione")
    run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = RGBColor.from_string(ISEO_RED)
    if spec.get("subtitle"):
        st = add_textbox(cover, Inches(0.6), Inches(3.3), W - Inches(1.2), Inches(0.8))
        r2 = st.text_frame.paragraphs[0].add_run()
        r2.text = spec["subtitle"]; r2.font.size = Pt(16)

    # slide di contenuto
    for sl in spec.get("slides", []):
        s = prs.slides.add_slide(layout)
        th = add_textbox(s, Inches(0.6), Inches(0.4), W - Inches(1.2), Inches(0.9))
        tr = th.text_frame.paragraphs[0].add_run()
        tr.text = sl.get("title", ""); tr.font.size = Pt(22); tr.font.bold = True
        tr.font.color.rgb = RGBColor.from_string(ISEO_RED)
        body = add_textbox(s, Inches(0.7), Inches(1.5), W - Inches(1.4), Inches(3.4))
        tf = body.text_frame
        first = True
        for b in sl.get("bullets", []):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run(); run.text = "•  " + str(b); run.font.size = Pt(16)
            p.space_after = Pt(6)

    name = _safe_name(spec.get("title", "presentazione"), "pptx")
    path = OUT_DIR / (next(tempfile._get_candidate_names()) + ".pptx")
    prs.save(str(path))
    return str(path), name


def gen_xlsx(spec: dict) -> tuple[str, str]:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    header_fill = PatternFill("solid", fgColor=ISEO_RED)
    header_font = Font(color="FFFFFF", bold=True)

    sheets = spec.get("sheets") or [{"name": "Foglio1", "columns": [], "rows": []}]
    for sh in sheets:
        ws = wb.create_sheet(title=(sh.get("name") or "Foglio")[:31])
        cols = sh.get("columns", [])
        for j, c in enumerate(cols, 1):
            cell = ws.cell(row=1, column=j, value=str(c))
            cell.fill = header_fill; cell.font = header_font
            cell.alignment = Alignment(horizontal="center")
        rows = sh.get("rows", [])
        for i, row in enumerate(rows, 2):
            for j, val in enumerate(row, 1):
                ws.cell(row=i, column=j, value=val)
        # totali con formula SOMMA per le colonne indicate
        totals = sh.get("total_columns", [])
        if rows and totals:
            tr = len(rows) + 2
            ws.cell(row=tr, column=1, value="TOTALE").font = Font(bold=True)
            for idx in totals:
                col = idx + 1
                if col <= len(cols):
                    letter = get_column_letter(col)
                    ws.cell(row=tr, column=col,
                            value=f"=SUM({letter}2:{letter}{len(rows)+1})").font = Font(bold=True)
        # larghezza colonne
        for j, c in enumerate(cols, 1):
            width = max(12, min(40, len(str(c)) + 4))
            ws.column_dimensions[get_column_letter(j)].width = width

    name = _safe_name(spec.get("title", "dati"), "xlsx")
    path = OUT_DIR / (next(tempfile._get_candidate_names()) + ".xlsx")
    wb.save(str(path))
    return str(path), name


def _docx_to_pdf(docx_path: str) -> str | None:
    """Converte docx→pdf con LibreOffice headless, se disponibile."""
    import shutil
    import subprocess
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None
    try:
        subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir",
                        str(OUT_DIR), docx_path], timeout=120,
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        pdf = Path(docx_path).with_suffix(".pdf")
        return str(pdf) if pdf.is_file() else None
    except Exception:
        return None


def _pdf_reportlab(spec: dict) -> tuple[str, str]:
    """Fallback PDF in stile ISEO (se LibreOffice non è disponibile)."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem

    name = _safe_name(spec.get("title", "documento"), "pdf")
    path = OUT_DIR / (next(tempfile._get_candidate_names()) + ".pdf")
    styles = getSampleStyleSheet()
    red = colors.HexColor("#" + ISEO_RED)
    h_title = ParagraphStyle("ISEOTitle", parent=styles["Title"], textColor=red, fontSize=22)
    h_sec = ParagraphStyle("ISEOSec", parent=styles["Heading1"], textColor=red, fontSize=14)
    body = styles["BodyText"]
    doc = SimpleDocTemplate(str(path), pagesize=A4, topMargin=2*cm, bottomMargin=2*cm)
    flow = [Paragraph(spec.get("title", "Documento"), h_title)]
    if spec.get("subtitle"):
        flow.append(Paragraph(spec["subtitle"], body))
    flow.append(Spacer(1, 0.4*cm))
    for sec in spec.get("sections", []):
        if sec.get("heading"):
            flow.append(Paragraph(sec["heading"], h_sec))
        for par in sec.get("paragraphs", []):
            flow.append(Paragraph(str(par), body)); flow.append(Spacer(1, 0.15*cm))
        bullets = sec.get("bullets", [])
        if bullets:
            flow.append(ListFlowable([ListItem(Paragraph(str(b), body)) for b in bullets],
                                     bulletType="bullet"))
        flow.append(Spacer(1, 0.25*cm))
    doc.build(flow)
    return str(path), name


def gen_pdf(spec: dict, template_path: str | None = None) -> tuple[str, str]:
    # preferito: Word sul template (ISEO o personale) -> PDF fedele
    try:
        docx_path, _ = gen_docx(spec, template_path=template_path)
        pdf = _docx_to_pdf(docx_path)
        if pdf:
            return pdf, _safe_name(spec.get("title", "documento"), "pdf")
    except Exception:
        pass
    # fallback: PDF in stile ISEO con reportlab
    return _pdf_reportlab(spec)


_BUILDERS = {"docx": gen_docx, "pptx": gen_pptx, "xlsx": gen_xlsx, "pdf": gen_pdf}


def generate(fmt: str, user_request: str, context: str, settings: dict,
             history_text: str = "", templates: dict | None = None,
             note_utente: str = "") -> tuple[str, str]:
    """Genera il file richiesto. Ritorna (path, filename_visibile).
    `templates` = {"docx": path|None, "pptx": path|None}: template PERSONALI
    dell'utente che, se presenti, BYPASSANO quelli di default (pdf usa il
    template Word). Un template invalido interrompe con errore parlante."""
    spec = build_spec(fmt, user_request, context, settings, history_text,
                      note_utente=note_utente)
    tpl = templates or {}
    if fmt in ("docx", "pdf"):
        return _BUILDERS[fmt](spec, template_path=tpl.get("docx"))
    if fmt == "pptx":
        return _BUILDERS[fmt](spec, template_path=tpl.get("pptx"))
    return _BUILDERS[fmt](spec)
