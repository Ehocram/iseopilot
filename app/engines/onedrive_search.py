#!/usr/bin/env python3
"""
OneDrive Search — Microsoft Graph API
Sviluppato da Marco Bonometti
Ricerca full-text nei file OneDrive tramite Microsoft Graph Search API.

OTTIMIZZAZIONI v2.1 (rispetto alla versione precedente)
-------------------------------------------------------
1) DOWNLOAD IN PARALLELO. Prima i contenuti dei file trovati venivano scaricati
   e ri-parsati uno per uno, in serie: con 5 risultati erano fino a 10 round-trip
   di rete sequenziali (metadati + download per ciascuno). Ora i download
   avvengono in parallelo (ThreadPoolExecutor): la latenza percepita scende
   tipicamente di 3-5×, a parità di risultati.
2) CACHE CONTENUTI. Lo stesso file, se ricompare in query successive, non viene
   ri-scaricato né ri-parsato: il testo estratto è messo in cache per
   (drive_id, item_id, data_modifica). Se il file cambia su OneDrive, la chiave
   cambia e viene riscaricato — niente risultati stantii.
3) LOG DI DEBUG DISATTIVABILE. Il log scriveva i NOMI dei file in chiaro su disco
   ad ogni fetch (I/O ad ogni chiamata + igiene del dato). Ora è disattivato per
   default e si abilita solo con la variabile d'ambiente CHAT_ASSISTANT_DEBUG=1.

L'interfaccia pubblica (OneDriveSearch.search, TokenManager, get_od) è invariata:
nessuna modifica necessaria in chat_assistant.py.
"""
import os
import json
import time
import datetime
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed

TOKEN_FILE = Path.home() / ".chat_assistant_od_token.json"

# Numero massimo di download paralleli (I/O-bound: si può alzare senza problemi)
_MAX_PARALLEL_FETCH = 5

# Cache testo estratto: chiave (drive_id, item_id, modified) -> testo. Cap semplice.
_CONTENT_CACHE = {}
_CONTENT_CACHE_MAX = 256


def _docx_all_text(doc) -> str:
    """Tutto il testo di un .docx già aperto: paragrafi + celle di tabella (anche
    annidate). python-docx con .paragraphs salta il testo dentro le TABELLE; i
    documenti ISEO sono a tabelle, quindi senza questo si estrarrebbe solo il
    titolo (era la causa dei contenuti OneDrive troncati)."""
    out = []
    def _walk(c):
        for p in getattr(c, "paragraphs", []):
            if p.text and p.text.strip():
                out.append(p.text)
        for t in getattr(c, "tables", []):
            for row in t.rows:
                for cell in row.cells:
                    _walk(cell)
    _walk(doc)
    return "\n".join(out)

# Log di debug solo se richiesto esplicitamente (default OFF)
_DEBUG = os.environ.get("CHAT_ASSISTANT_DEBUG", "") == "1"


def _odlog(msg: str):
    if not _DEBUG:
        return
    try:
        with open(Path.home() / "chat_assistant_debug.txt", "a", encoding="utf-8") as f:
            f.write(f"[ODFETCH] {msg}\n")
    except Exception:
        pass


# ── Token manager (invariato) ────────────────────────────
class TokenManager:
    """Gestisce access token e refresh token Microsoft OAuth2."""

    def __init__(self, client_id: str, tenant_id: str = "common"):
        self.client_id = client_id
        self.tenant_id = tenant_id
        self._token_data = None
        self._cancel_poll = False
        self._poll_in_progress = False
        self._load()

    def cancel_polling(self):
        self._cancel_poll = True

    def _load(self):
        try:
            if TOKEN_FILE.exists():
                self._token_data = json.loads(TOKEN_FILE.read_text(encoding="utf-8"))
        except Exception:
            self._token_data = None

    def _save(self, data: dict):
        try:
            self._token_data = data
            TOKEN_FILE.write_text(json.dumps(data, indent=2), encoding="utf-8")
        except Exception:
            pass

    def is_authenticated(self) -> bool:
        return bool(self._token_data and self._token_data.get("access_token"))

    def get_access_token(self) -> str:
        if not self._token_data:
            return ""
        expires_at   = self._token_data.get("expires_at", 0)
        access_token = self._token_data.get("access_token", "")
        if access_token and time.time() < expires_at - 300:
            return access_token
        refresh_token = self._token_data.get("refresh_token", "")
        if refresh_token:
            new_token = self._refresh(refresh_token)
            if new_token:
                return new_token
        return access_token

    def _refresh(self, refresh_token: str) -> str:
        try:
            import requests as req
            url = f"https://login.microsoftonline.com/{self.tenant_id}/oauth2/v2.0/token"
            r = req.post(url, data={
                "client_id":     self.client_id,
                "grant_type":    "refresh_token",
                "refresh_token": refresh_token,
                "scope":         "Files.Read.All offline_access",
            }, timeout=15)
            resp = r.json()
            if "access_token" in resp:
                resp["expires_at"] = time.time() + resp.get("expires_in", 3600)
                self._save(resp)
                return resp["access_token"]
        except Exception:
            pass
        return ""

    def start_device_flow(self) -> dict:
        try:
            import requests as req
            url = f"https://login.microsoftonline.com/{self.tenant_id}/oauth2/v2.0/devicecode"
            r = req.post(url, data={
                "client_id": self.client_id,
                "scope":     "Files.Read.All offline_access",
            }, timeout=15)
            return r.json()
        except Exception as e:
            return {"error": str(e)}

    def poll_device_flow(self, device_code: str, interval: int = 5) -> tuple:
        if self._poll_in_progress:
            return False, "Polling gia in corso — attendi o riavvia l'app"
        self._poll_in_progress = True
        self._cancel_poll = False
        try:
            import requests as req
            url = f"https://login.microsoftonline.com/{self.tenant_id}/oauth2/v2.0/token"
            poll_interval = max(interval, 3)
            for attempt in range(80):
                slept = 0.0
                while slept < poll_interval:
                    if self._cancel_poll:
                        return False, "Polling annullato dall'utente"
                    time.sleep(0.5)
                    slept += 0.5
                if self._cancel_poll:
                    return False, "Polling annullato dall'utente"
                try:
                    r = req.post(url, data={
                        "client_id":   self.client_id,
                        "grant_type":  "urn:ietf:params:oauth:grant-type:device_code",
                        "device_code": device_code,
                    }, timeout=10)
                    resp = r.json()
                    if "access_token" in resp:
                        resp["expires_at"] = time.time() + resp.get("expires_in", 3600)
                        self._save(resp)
                        return True, "Autenticato con successo"
                    err = resp.get("error", "")
                    if err == "authorization_pending":
                        continue
                    elif err == "slow_down":
                        poll_interval += 5
                        continue
                    elif err == "expired_token":
                        return False, "Codice scaduto — riprova il login"
                    elif err == "access_denied":
                        return False, "Accesso negato dall'utente"
                    else:
                        return False, resp.get("error_description", f"Errore: {err}")
                except Exception:
                    continue
            return False, "Timeout — il codice e scaduto. Riprova il login."
        except Exception as e:
            return False, str(e)
        finally:
            self._poll_in_progress = False
            self._cancel_poll = False

    def logout(self):
        self.cancel_polling()
        self._token_data = None
        try:
            if TOKEN_FILE.exists():
                TOKEN_FILE.unlink()
        except Exception:
            pass


# ── Graph Search ─────────────────────────────────────────
def _build_query(text: str) -> str:
    """Estrae parole chiave per OneDrive mappando i CONCETTI a TERMINI
    DOCUMENTALI (porting fedele di _extract_keywords dell'app desktop).

    Esempio chiave: "quanti anni ha Marco Bonometti?" non viene cercato così
    com'è (Graph cerca i termini), ma tradotto in "anagrafica bonometti marco":
    il concetto 'anni' mappa al documento 'anagrafica', che è il file giusto.
    """
    import re
    tl = (text or "").lower()

    # Mappa concettuale — concetto della domanda -> termine presente nei documenti.
    concept_map = [
        (['anni', 'anno', 'nato', 'nata', 'nascita', 'eta', 'età', 'compleanno',
          'quanti anni', 'data nascita', 'quand',
          'how old', 'years old', 'born', 'birth', 'age of'], 'anagrafica'),
        (['risiede', 'abita', 'residente', 'residenza', 'indirizzo', 'dove vive',
          'città', 'citta', 'comune', 'domicilio',
          'address', 'lives in', 'resides'], 'anagrafica'),
        (['codice fiscale', 'partita iva', 'piva', 'cf ',
          'tax code', 'vat number', 'fiscal code'], 'anagrafica'),
        (['contratto', 'assunto', 'assunzione', 'stipendio', 'salario', 'ral',
          'ferie', 'permessi', 'ccnl', 'busta paga',
          'contract', 'salary', 'hired', 'hiring', 'payslip', 'vacation'], 'contratto'),
        (['fattura', 'invoice', 'pagamento', 'importo', 'bonifico',
          'payment', 'bank transfer', 'amount due'], 'fattura'),
        (['vulnerabilit', 'cve', 'exploit', 'patch', 'pentest', 'remediation',
          'vulnerability'], 'vulnerability'),
        (['riunione', 'meeting', 'verbale', 'minuti riunione',
          'meeting minutes', 'minutes of'], 'verbale'),
        (['report', 'relazione', 'analisi', 'findings', 'risultati',
          'analysis', 'results'], 'report'),
    ]

    extra = []
    for triggers, keyword in concept_map:
        if any(t in tl for t in triggers):
            extra.append(keyword)
            if len(extra) >= 2:
                break

    # Stopword — esclude le parole funzionali (non i termini della concept_map).
    stopwords = {
        'ciao', 'il', 'lo', 'la', 'gli', 'le', 'un', 'uno', 'una', 'di', 'da', 'in',
        'con', 'su', 'per', 'tra', 'fra', 'e', 'o', 'ma', 'se', 'non', 'che', 'chi',
        'come', 'quando', 'dove', 'cosa', 'qual', 'quali', 'quanto', 'quanti', 'quante',
        'mi', 'ti', 'ci', 'vi', 'si', 'ne', 'ho', 'ha', 'hai', 'hanno', 'sono', 'sei',
        'siamo', 'siete', 'era', 'del', 'della', 'dei', 'degli', 'delle', 'al', 'alla',
        'ai', 'agli', 'alle', 'nel', 'nella', 'nei', 'negli', 'nelle', 'sul', 'sulla',
        'sui', 'sugli', 'sulle', 'dal', 'dalla', 'dai', 'the', 'an', 'is', 'are', 'was',
        'were', 'been', 'have', 'has', 'had', 'did', 'will', 'would', 'could', 'should',
        'may', 'might', 'shall', 'can', 'dammi', 'dimmi', 'parlami', 'raccontami',
        'mostrami', 'dicci', 'descrivimi', 'quale', 'questo', 'questa', 'questi',
        'queste', 'molto', 'poco', 'tanto', 'fatto', 'fatta', 'fatti', 'fatte',
        'avere', 'essere', 'stare',
    }

    words = re.findall(r'[a-zA-ZàèéìòùÀÈÉÌÒÙ]{4,}', tl)
    concept_kw = {kw for _, kw in concept_map}
    content_words = [w for w in words if w not in stopwords and w not in concept_kw]
    content_words = sorted(set(content_words), key=len, reverse=True)[:2]

    result = extra + [w for w in content_words if w not in extra]
    return ' '.join(result[:3]) if result else (text or "")[:50]


# Estensioni senza testo estraibile: inquinano i risultati documentali.
_BIN_EXT = (".zip", ".rar", ".7z", ".gz", ".tar", ".exe", ".dmg",
            ".iso", ".msi", ".pkg", ".bin", ".apk", ".img",
            # immagini: nessun testo estraibile, non competono nei risultati
            ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg", ".heic", ".webp")


def _is_binary_name(name: str) -> bool:
    return str(name or "").lower().endswith(_BIN_EXT)


class OneDriveSearch:
    """Ricerca full-text in OneDrive tramite Microsoft Graph Search API."""

    GRAPH_SEARCH = "https://graph.microsoft.com/v1.0/search/query"
    GRAPH_ME     = "https://graph.microsoft.com/v1.0/me"

    def __init__(self, cfg: dict):
        self.cfg = cfg
        self.tm  = TokenManager(
            client_id=cfg.get("od_client_id", ""),
            tenant_id=cfg.get("od_tenant_id", "common"),
        )

    def is_configured(self) -> bool:
        return bool(self.cfg.get("od_client_id", "").strip())

    def is_authenticated(self) -> bool:
        return self.tm.is_authenticated() and bool(self.tm.get_access_token())

    def get_user_info(self) -> dict:
        try:
            import requests as req
            token = self.tm.get_access_token()
            if not token:
                return {}
            r = req.get(self.GRAPH_ME,
                headers={"Authorization": f"Bearer {token}"}, timeout=10)
            return r.json()
        except Exception:
            return {}

    def search(self, query: str, max_results: int = 5) -> str:
        """
        Cerca in OneDrive e ritorna testo formattato con i risultati.
        Ritorna stringa vuota se nessun risultato o errore.
        I contenuti dei file trovati vengono scaricati IN PARALLELO.
        """
        token = self.tm.get_access_token()
        self.last_links = []   # link strutturati (nome, url) dell'ultima ricerca
        if not token:
            return ""
        # Domanda -> parole chiave: Graph cerca i termini, non interpreta la frase.
        query_string = _build_query(query)
        _odlog(f"OneDrive query: '{query}' -> '{query_string}'")
        try:
            import requests as req
            payload = {
                "requests": [{
                    "entityTypes": ["driveItem"],
                    "query": {"queryString": query_string},
                    "from": 0,
                    "size": max_results,
                    "fields": ["id", "name", "webUrl", "lastModifiedDateTime",
                               "file", "parentReference", "summary"],
                }]
            }
            r = req.post(self.GRAPH_SEARCH,
                json=payload,
                headers={"Authorization": f"Bearer {token}"},
                timeout=20)
            if r.status_code == 401:
                new_token = self.tm._refresh(self.tm._token_data.get("refresh_token", ""))
                if new_token:
                    token = new_token
                    r = req.post(self.GRAPH_SEARCH,
                        json=payload,
                        headers={"Authorization": f"Bearer {new_token}"},
                        timeout=20)
            data = r.json()

            hits = (data.get("value", [{}])[0]
                        .get("hitsContainers", [{}])[0]
                        .get("hits", []))
            if not hits:
                return ""

            # ── 1) Estrai i metadati di ogni hit (operazione locale, niente rete) ──
            metas = []
            for hit in hits:
                res     = hit.get("resource", {})
                summary = hit.get("summary", "").strip()

                name = (res.get("name")
                        or res.get("displayName")
                        or res.get("listItem", {}).get("fields", {}).get("FileLeafRef")
                        or hit.get("hitId", "?"))
                url = (res.get("webUrl")
                       or res.get("listItem", {}).get("webUrl", ""))
                fsi = res.get("fileSystemInfo", {})
                modified = (fsi.get("lastModifiedDateTime", "")
                            or res.get("lastModifiedDateTime", ""))
                parent = res.get("parentReference", {})
                folder = parent.get("path", "").split("root:")[-1] if parent else ""
                drive_id = parent.get("driveId", "")
                item_id  = (res.get("id", "")
                            or res.get("listItem", {}).get("id", "")
                            or hit.get("hitId", ""))
                metas.append({
                    "name": name, "url": url, "modified": modified,
                    "folder": folder, "drive_id": drive_id, "item_id": item_id,
                    "summary": summary,
                })

            # Filtro qualità: archivi e binari (zip, exe, dmg…) non hanno testo
            # estraibile e tolgono posti ai documenti veri. Vengono scartati,
            # ma SOLO se resta almeno un risultato documentale (mai zero
            # risultati per colpa del filtro).
            textual = [m for m in metas if not _is_binary_name(m.get("name"))]
            if textual and len(textual) < len(metas):
                _odlog(f"filtrati {len(metas) - len(textual)} file binari/archivio dai risultati")
                metas = textual

            # ── 2) Scarica i contenuti IN PARALLELO (I/O-bound) ──
            contents = [""] * len(metas)
            workers = min(_MAX_PARALLEL_FETCH, len(metas)) or 1
            with ThreadPoolExecutor(max_workers=workers) as ex:
                fut_to_i = {
                    ex.submit(self._fetch_file_content, token,
                              m["drive_id"], m["item_id"], m["name"], m["modified"]): i
                    for i, m in enumerate(metas)
                }
                for fut in as_completed(fut_to_i):
                    i = fut_to_i[fut]
                    try:
                        contents[i] = fut.result() or ""
                    except Exception:
                        contents[i] = ""

            # ── 3) Assembla i risultati nell'ordine originale (rilevanza Graph) ──
            parts = []
            for m, content in zip(metas, contents):
                chunk = f"[OneDrive: {m['name']}]"
                if m["folder"]:
                    chunk += f" (cartella: {m['folder']})"
                if m["modified"]:
                    chunk += f" — modificato: {m['modified'][:10]}"
                if content:
                    chunk += f"\n{content}"
                elif m["summary"]:
                    chunk += f"\n{m['summary']}"
                if m["url"]:
                    chunk += f"\nLink: {m['url']}"
                parts.append(chunk)

            # Link strutturati (nome, url) per chi consuma i risultati: evita di dover
            # ri-parsare il testo con regex (fragile se un file non ha webUrl).
            self.last_links = [(m["name"], m["url"]) for m in metas if m.get("url")]
            return "\n\n".join(parts)

        except Exception as e:
            _odlog(f"search EXCEPTION: {type(e).__name__}: {e}")
            return ""

    def _fetch_file_content(self, token: str, drive_id: str, item_id: str,
                            name: str, modified: str = "", max_chars: int = 6000) -> str:
        """Scarica e legge il contenuto di un file da OneDrive via Graph API.
        Usa una cache per (drive_id, item_id, modified): se il file non è cambiato,
        non viene ri-scaricato né ri-parsato."""
        _odlog(f"start: name={name!r} drive_id={drive_id!r} item_id={item_id!r}")
        if not drive_id or not item_id:
            _odlog("ABORT: drive_id o item_id mancanti")
            return ""

        ext = name.lower().rsplit(".", 1)[-1] if "." in name else ""
        if ext not in ("pdf", "docx", "doc", "txt", "md", "csv", "xlsx", "xls", "pptx"):
            return ""  # salta file non leggibili (immagini, zip, ecc.)

        # ── cache ──
        ck = (drive_id, item_id, modified)
        if ck in _CONTENT_CACHE:
            _odlog("cache HIT")
            return _CONTENT_CACHE[ck]

        try:
            import requests as req, io
            r = req.get(
                f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{item_id}"
                f"?select=@microsoft.graph.downloadUrl,name",
                headers={"Authorization": f"Bearer {token}"}, timeout=15)
            _odlog(f"metadata HTTP {r.status_code}")
            if r.status_code != 200:
                _odlog(f"metadata body: {r.text[:200]}")
                return ""
            download_url = r.json().get("@microsoft.graph.downloadUrl", "")
            if not download_url:
                _odlog("downloadUrl mancante nella risposta")
                return ""

            r2 = req.get(download_url, timeout=30, stream=True)
            _odlog(f"download HTTP {r2.status_code}")
            if r2.status_code != 200:
                return ""
            raw = b""
            for chunk in r2.iter_content(chunk_size=65536):
                raw += chunk
                if len(raw) > 5 * 1024 * 1024:   # cap 5 MB
                    break

            text = self._extract(raw, ext)
            _odlog(f"text estratto: {len(text)} chars (ext={ext})")
            result = text.strip()[:max_chars] if text.strip() else ""

            # memorizza in cache (con cap FIFO semplice)
            if len(_CONTENT_CACHE) >= _CONTENT_CACHE_MAX:
                try:
                    _CONTENT_CACHE.pop(next(iter(_CONTENT_CACHE)))
                except Exception:
                    _CONTENT_CACHE.clear()
            _CONTENT_CACHE[ck] = result
            return result
        except Exception as e:
            _odlog(f"EXCEPTION: {type(e).__name__}: {e}")
            return ""

    @staticmethod
    def _extract(raw: bytes, ext: str) -> str:
        import io
        text = ""
        try:
            if ext == "pdf":
                import fitz
                doc = fitz.open(stream=raw, filetype="pdf")
                text = "\n".join(page.get_text() for page in doc)
            elif ext in ("docx", "doc"):
                from docx import Document
                doc = Document(io.BytesIO(raw))
                text = _docx_all_text(doc)
            elif ext in ("txt", "md", "csv"):
                text = raw.decode("utf-8", errors="ignore")
            elif ext in ("xlsx", "xls"):
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
                rows = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        r_str = "\t".join(str(c) if c is not None else "" for c in row)
                        if r_str.strip():
                            rows.append(r_str)
                text = "\n".join(rows)
            elif ext == "pptx":
                from pptx import Presentation
                prs = Presentation(io.BytesIO(raw))
                slides = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slides.append(shape.text.strip())
                text = "\n".join(slides)
        except Exception:
            return ""
        return text


# ── Istanza globale ──────────────────────────────────────
_od_instance = None

def get_od(cfg: dict) -> OneDriveSearch:
    global _od_instance
    if _od_instance is None or _od_instance.cfg.get("od_client_id") != cfg.get("od_client_id"):
        _od_instance = OneDriveSearch(cfg)
    else:
        _od_instance.cfg = cfg
    return _od_instance
