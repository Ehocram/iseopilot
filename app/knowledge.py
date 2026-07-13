"""
knowledge.py — Livello conoscenza per la web app.

Riusa i moduli desktop senza modificarli:
  * engines/vector_db.py     -> ChromaDB. Collezione = dipartimento (kb_<area>).
  * engines/folder_index.py  -> ricerca FTS5 su cartelle locali o di rete.

Isolamento:
  * Ogni dipartimento ha la PROPRIA collezione ChromaDB. Un utente carica e
    interroga solo la collezione del suo dipartimento.
  * I dati ChromaDB e gli indici cartelle vivono sotto APP_DATA_DIR (volume),
    non nella home dell'utente.

Note dipendenze:
  * La ricerca CARTELLE su .txt/.md/.csv/.json funziona con la sola libreria
    standard (FTS5). PDF/DOCX/XLSX richiedono i pacchetti in
    requirements-connectors.txt (PyMuPDF/python-docx/openpyxl).
  * ChromaDB richiede `pip install -r requirements-connectors.txt`. Se assente,
    le funzioni KB ritornano un errore esplicito (nessun fallimento silenzioso).
"""
from __future__ import annotations

import os
import tempfile
import threading
from pathlib import Path

from . import store
from .engines import vector_db, folder_index

DATA_DIR = Path(os.environ.get("APP_DATA_DIR", "/data"))
CHROMA_DIR = DATA_DIR / "vectordb"
FOLDER_INDEX_DIR = DATA_DIR / "folder_index"
CHROMA_DIR.mkdir(parents=True, exist_ok=True)
FOLDER_INDEX_DIR.mkdir(parents=True, exist_ok=True)

# Reindirizza gli indici cartella sotto il volume dati (non nella home).
folder_index.INDEX_DIR = FOLDER_INDEX_DIR

# Estensioni accettate per l'upload in ChromaDB.
ALLOWED_UPLOAD_EXT = {".txt", ".md", ".csv", ".json", ".xml", ".py",
                      ".pdf", ".docx", ".xlsx", ".xls"}
# Gli allegati di chat accettano anche le presentazioni.
ALLOWED_ATTACH_EXT = ALLOWED_UPLOAD_EXT | {".pptx", ".ppt", ".rtf"}


# ── Configurazione ChromaDB per dipartimento ────────────────
def vector_cfg(dept: str) -> dict:
    """cfg per VectorDBManager: collezione = dipartimento; infra da admin."""
    return {
        "db_mode": store.get_setting("kb_mode", "local"),
        "db_embedding_model": store.get_setting("kb_embedding_model", "all-minilm-l6-v2"),
        "db_collection": store.collection_for_department(dept),
        "db_path": str(CHROMA_DIR),
        "db_host": store.get_setting("kb_host", "localhost"),
        "db_port": store.get_setting("kb_port", "8000"),
        "db_api_key": store.get_setting("kb_api_key", ""),
    }


def kb_available() -> bool:
    try:
        import chromadb  # noqa: F401
        return True
    except Exception:
        return False


# ── Estrazione testo dai file caricati ──────────────────────
def extract_text(filename: str, raw: bytes) -> str:
    """Estrae il testo da un file caricato (riusa l'estrattore desktop)."""
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_UPLOAD_EXT:
        return ""
    suffix = ext or ".bin"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(raw)
        tmp_path = Path(tmp.name)
    try:
        return folder_index._extract_text(tmp_path, ext) or ""
    finally:
        try:
            tmp_path.unlink()
        except Exception:
            pass


def _extract_pptx_text(raw: bytes) -> str:
    """Testo di una presentazione: titoli, testo delle forme, note relatore."""
    try:
        from pptx import Presentation
        import io
        prs = Presentation(io.BytesIO(raw))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"-- Slide {i} --")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in p.runs).strip()
                        if t:
                            out.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        out.append(" | ".join(cells))
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note = slide.notes_slide.notes_text_frame.text.strip()
                if note:
                    out.append("[Note: " + note + "]")
        return "\n".join(out)
    except Exception:
        return ""


# Tetto di sicurezza per gli ALLEGATI di chat: molto ampio (2M caratteri).
# NB: l'estrattore delle cartelle ha un tetto a 400k pensato per l'INDICE FTS;
# per gli allegati serve fedeltà totale ("quanti X ci sono?" deve contare
# TUTTE le righe), quindi i tabellari hanno un percorso dedicato senza quel tetto.
ATTACHMENT_MAX_CHARS = 2_000_000


def _extract_xlsx_full(raw: bytes) -> str:
    """Excel per allegati: TUTTI i fogli, TUTTE le righe (riga per riga,
    tab-separate), senza il tetto da indicizzazione."""
    import io
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
    lines, tot = [], 0
    for ws in wb.worksheets:
        lines.append(f"[Foglio: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            ln = "\t".join(str(c) if c is not None else "" for c in row)
            lines.append(ln)
            tot += len(ln) + 1
            if tot > ATTACHMENT_MAX_CHARS:
                lines.append("[…estrazione interrotta al tetto di sicurezza…]")
                return "\n".join(lines)
    return "\n".join(lines)


def extract_attachment_text(filename: str, raw: bytes) -> str:
    """Estrae il testo da un allegato di chat. Come extract_text ma include
    anche le presentazioni PowerPoint e, per i tabellari, l'estrazione
    integrale (fedeltà da prompt: i conteggi devono tornare)."""
    ext = Path(filename).suffix.lower()
    if ext in {".pptx", ".ppt"}:
        return _extract_pptx_text(raw)
    if ext in {".xlsx", ".xlsm"}:
        try:
            return _extract_xlsx_full(raw)
        except Exception:
            pass  # ripiego sul percorso standard
    if ext == ".rtf":
        # I Mac producono RTF di default (TextEdit): rifiutarlo generava solo
        # chip d'errore e confusione. striprtf è tollerante sui file reali.
        try:
            from striprtf.striprtf import rtf_to_text
            return rtf_to_text(raw.decode("utf-8", errors="replace"))[:ATTACHMENT_MAX_CHARS]
        except Exception:
            pass
    if ext == ".csv":
        try:
            return raw.decode("utf-8", errors="replace")[:ATTACHMENT_MAX_CHARS]
        except Exception:
            pass
    if ext in ALLOWED_UPLOAD_EXT:
        return extract_text(filename, raw)
    return ""


# ── Operazioni KB (ChromaDB) per dipartimento ───────────────
def kb_ingest(dept: str, filename: str, text: str) -> tuple[bool, str]:
    if not kb_available():
        return False, "ChromaDB non installato sul server (pip install -r requirements-connectors.txt)."
    if not text.strip():
        return False, f"Nessun testo estraibile da '{filename}'."
    vdb = vector_db.get_vdb(vector_cfg(dept))
    n, msg = vdb.add_document(text, filename)
    return (n > 0), msg


def kb_list(dept: str) -> list:
    if not kb_available():
        return []
    try:
        return vector_db.get_vdb(vector_cfg(dept)).list_documents()
    except Exception:
        return []


def kb_count(dept: str) -> int:
    if not kb_available():
        return 0
    try:
        return vector_db.get_vdb(vector_cfg(dept)).get_count()
    except Exception:
        return 0


def kb_delete(dept: str, filename: str) -> tuple[bool, str]:
    if not kb_available():
        return False, "ChromaDB non installato."
    try:
        ok, msg = vector_db.get_vdb(vector_cfg(dept)).delete_document(filename)
        return bool(ok), msg
    except Exception as e:
        return False, f"Errore eliminazione: {e}"


def kb_search(dept: str, query: str, n: int = 4) -> tuple[str, list]:
    """Ricerca nella Conoscenza con GARANZIA DI COPERTURA per le domande di
    enumerazione (caso Carlos: 'quali certificazioni ISO abbiamo?' pescava 2
    documenti su 3). Tre mosse: (1) i documenti il cui NOME matcha i termini
    della domanda entrano TUTTI, col loro miglior passaggio; (2) il resto del
    budget è semantico ma DIVERSIFICATO per fonte (max 2 passaggi a documento);
    (3) la copertura è DICHIARATA al modello — contratto dei conteggi, come
    per gli allegati. In caso d'errore, ripiego sul comportamento storico."""
    if not kb_available():
        return "", []
    try:
        vdb = vector_db.get_vdb(vector_cfg(dept))
        try:
            import re as _re
            _stop = {"che", "come", "quali", "quale", "quante", "quanti", "una",
                     "uno", "gli", "del", "della", "delle", "dei", "per", "con",
                     "abbiamo", "sono", "nella", "nostra", "nostre", "the", "and",
                     "have", "does", "what", "which", "our"}
            terms = [w for w in _re.findall(r"[a-zà-ÿ0-9]{2,}", (query or "").lower())
                     if w not in _stop]
            nomi = [d.get("name", "") for d in vdb.list_documents()]
            scored = []
            for nome in nomi:
                nl = nome.lower()
                sc = sum(1 for t in terms if t in nl)
                if sc > 0:
                    scored.append((-sc, nome))
            scored.sort()
            corrispondenti = [nome for _s, nome in scored[:8]]

            parti, fonti, usati = [], [], set()
            budget = 5800
            # 1) copertura garantita: un passaggio da OGNI documento col nome
            for nome in corrispondenti:
                chunk = vdb.search_in_document(query, nome, n_results=1)
                if chunk.strip():
                    blocco = f"[Fonte: {nome}]\n{chunk[:900]}"
                    if sum(len(x) for x in parti) + len(blocco) > budget:
                        break
                    parti.append(blocco)
                    fonti.append((nome, ""))
                    usati.add(nome)
            # 2) semantico diversificato per fonte (max 2 passaggi a documento)
            per_fonte = {}
            for chunk, src in vdb.search_raw(query, n_results=max(n * 2, 8)):
                if per_fonte.get(src, 0) >= 2:
                    continue
                blocco = f"[Fonte: {src}]\n{chunk[:900]}"
                if sum(len(x) for x in parti) + len(blocco) > budget:
                    break
                parti.append(blocco)
                per_fonte[src] = per_fonte.get(src, 0) + 1
                if src not in usati:
                    fonti.append((src, ""))
                    usati.add(src)
            if not parti:
                return "", []
            testa = ""
            if corrispondenti:
                inclusi = [nome for nome in corrispondenti if nome in usati]
                testa = ("[DOCUMENTI della Conoscenza corrispondenti alla domanda "
                         f"per NOME: {'; '.join(corrispondenti)} — "
                         f"{'TUTTI inclusi qui sotto' if len(inclusi) == len(corrispondenti) else f'inclusi {len(inclusi)} su {len(corrispondenti)}'}. "
                         "Per domande di elenco/conteggio usa QUESTA lista.]\n\n")
            return testa + "\n\n".join(parti), fonti
        except Exception:
            return vdb.search(query, n_results=n)  # comportamento storico
    except Exception:
        return "", []


# ── Ricerca cartelle (locali / di rete) ─────────────────────
def folder_reindex(path: str) -> tuple[bool, str]:
    p = (path or "").strip()
    if not p:
        return False, "Nessuna cartella configurata."
    if not Path(p).exists():
        return False, f"Percorso non raggiungibile dal server: {p}"
    try:
        idx = folder_index.get_index(p)
        res = idx.update()
        n = idx.count()
        return True, f"Indice aggiornato: {n} passaggi indicizzati."
    except Exception as e:
        return False, f"Errore indicizzazione: {e}"


def folder_count(path: str) -> int:
    p = (path or "").strip()
    if not p or not Path(p).exists():
        return 0
    try:
        res = folder_index.get_index(p).count()
        # FolderIndex.count() -> (num_file, num_passaggi); ci interessano i passaggi.
        return res[1] if isinstance(res, (tuple, list)) else int(res)
    except Exception:
        return 0


def folder_search(path: str, query: str, n: int = 4, rerank: bool = False) -> tuple[str, list]:
    p = (path or "").strip()
    if not p or not Path(p).exists():
        return "", []
    try:
        return folder_index.search_folder(p, query, top_k=n, rerank=rerank)
    except Exception:
        return "", []


# ── Cartelle del dipartimento (più percorsi) ────────────────
def dept_folders_reindex(dept: str) -> tuple[bool, str]:
    paths = store.department_folders(dept)
    if not paths:
        return False, "Nessuna cartella configurata per il dipartimento."
    done, msgs = 0, []
    for p in paths:
        ok, msg = folder_reindex(p)
        if ok:
            done += 1
        else:
            msgs.append(msg)
    if done == 0:
        return False, " · ".join(msgs) or "Indicizzazione non riuscita."
    tail = (" (" + " · ".join(msgs) + ")") if msgs else ""
    return True, f"{done}/{len(paths)} cartelle indicizzate.{tail}"


def dept_folders_count(dept: str) -> int:
    return sum(folder_count(p) for p in store.department_folders(dept))


# ── Migrazione modello di embedding (re-indicizzazione KB) ──
def _embedding_ok(model_alias: str) -> tuple[bool, str]:
    """Verifica che il modello di embedding sia caricabile e funzionante PRIMA
    di toccare i dati (al primo uso viene scaricato: serve accesso di rete)."""
    try:
        from .engines.vector_db import resolve_embedding_model
        from chromadb.utils import embedding_functions
        ef = embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name=resolve_embedding_model(model_alias))
        vec = ef(["testo di prova"])
        return (bool(vec) and len(vec) == 1), "ok"
    except Exception as e:
        return False, str(e)[:300]


def kb_reembed_all(new_model: str = "multilingual") -> tuple[int, list[str]]:
    """Re-indicizza tutte le collezioni KB col nuovo modello di embedding.

    Necessario quando si cambia modello: i vettori vecchi e le query nuove
    vivrebbero in spazi diversi (silenziosamente, perché le dimensioni
    coincidono) e la ricerca degraderebbe. Ricostruisce ogni documento dai
    chunk già salvati in Chroma (metadata 'source') e lo re-ingerisce.
    SICUREZZA: il modello nuovo è verificato PRIMA di cancellare qualsiasi
    collezione; i dati sono letti in memoria prima della cancellazione."""
    if not kb_available():
        return 0, ["ChromaDB non installato sul server."]
    ok, msg = _embedding_ok(new_model)
    if not ok:
        raise RuntimeError(
            f"Modello '{new_model}' non caricabile ({msg}). Nessun dato toccato. "
            "Verifica che il server possa raggiungere huggingface.co per il primo download.")
    import chromadb
    client = chromadb.PersistentClient(path=str(CHROMA_DIR))
    store.set_setting("kb_embedding_model", new_model)
    done, esiti = 0, []
    for dept in store.list_departments():
        cname = store.collection_for_department(dept)
        try:
            col = client.get_collection(cname)
        except Exception:
            continue  # nessuna KB per questo reparto
        try:
            data = col.get(include=["documents", "metadatas"])
        except Exception as e:
            esiti.append(f"{dept}: ERRORE lettura ({str(e)[:80]})")
            continue
        by_file: dict[str, list] = {}
        for doc, meta in zip(data.get("documents") or [], data.get("metadatas") or []):
            src = (meta or {}).get("source", "documento")
            by_file.setdefault(src, []).append(((meta or {}).get("chunk", 0), doc or ""))
        if not by_file:
            continue
        # dati in memoria: ora si può ricreare la collezione col nuovo modello
        try:
            client.delete_collection(cname)
        except Exception:
            pass
        ok_files = 0
        for fname, chunks in by_file.items():
            chunks.sort(key=lambda t: t[0])
            text = "\n\n".join(c for _, c in chunks if c.strip())
            good, _m = kb_ingest(dept, fname, text)  # usa il nuovo modello dal setting
            if good:
                ok_files += 1
        done += 1
        esiti.append(f"{dept}: {ok_files}/{len(by_file)} file re-indicizzati")
    return done, esiti


def kb_reembed_async(new_model: str = "multilingual") -> None:
    """Lancia la re-indicizzazione in background e traccia l'esito in un
    setting leggibile dalla pagina admin (fail loudly, mai in silenzio)."""
    def _job():
        store.set_setting("kb_reembed_status", "in corso… (al primo avvio scarica il modello)")
        try:
            n, esiti = kb_reembed_all(new_model)
            store.set_setting("kb_reembed_status",
                              f"completato ({n} reparti) — " + "; ".join(esiti[:8]) if esiti
                              else "completato: nessuna KB da re-indicizzare")
        except Exception as e:
            store.set_setting("kb_reembed_status", "ERRORE: " + str(e)[:400])
    threading.Thread(target=_job, daemon=True).start()


# ── Indicizzazione automatica (come l'app desktop: all'avvio + on demand) ──
def reindex_all() -> None:
    """Indicizza (incrementale) tutte le cartelle di tutti i dipartimenti.
    Le share di rete grandi possono richiedere tempo: va eseguita in background.
    L'update di FolderIndex è incrementale (mtime+size), quindi ai riavvii
    successivi rilegge solo i file nuovi o modificati."""
    seen = set()
    for dept in store.list_departments():
        for p in store.department_folders(dept):
            if p in seen:
                continue
            seen.add(p)
            try:
                folder_reindex(p)
            except Exception:
                pass


def reindex_folder_async(path: str) -> None:
    """Indicizza una singola cartella in un thread separato (non blocca la risposta)."""
    threading.Thread(target=folder_reindex, args=(path,), daemon=True).start()


def reindex_all_async() -> None:
    """Avvia in background l'indicizzazione di tutte le cartelle configurate."""
    threading.Thread(target=reindex_all, daemon=True).start()


# ── Recupero combinato per la chat (RAG) ────────────────────
def enrich_query(query: str, prev_user_query: str = "") -> str:
    """Arricchisce la domanda per la ricerca documentale (KB, cartelle, connettori).

    1) Follow-up: se la domanda è breve o anaforica ("e per il 2025?"), eredita
       i termini di contenuto del turno utente precedente, altrimenti la ricerca
       non ha soggetto e non trova nulla.
    2) Concept map (stessa dell'app desktop): traduce il concetto della domanda
       nel termine documentale ("quanti anni ha X" -> aggiunge "anagrafica")."""
    import re as _re
    q = (query or "").strip()
    # Parole di contenuto = tolte stopword e pronomi anaforici ("quelli",
    # "questo"...): i pronomi sono il segnale del follow-up, non un soggetto.
    _anaphoric = {"quello", "quella", "quelli", "quelle", "questo", "questa",
                  "questi", "queste", "stesso", "stessa", "stessi", "stesse",
                  "altro", "altra", "altri", "altre", "invece", "anche",
                  "quindi", "allora", "prima", "dopo", "ancora", "pure",
                  "that", "this", "those", "these", "same", "other", "ones"}
    try:
        from .engines.folder_index import _STOP as _stop
    except Exception:
        _stop = set()
    content = [w for w in _re.findall(r"[A-Za-zÀ-ÿ0-9]{3,}", q.lower())
               if w not in _stop and w not in _anaphoric]
    if prev_user_query and len(content) < 2:
        q = (prev_user_query.strip() + " " + q).strip()
    try:
        from .engines.onedrive_search import _build_query as _kw
        extra = _kw(q)
        if extra and extra.strip():
            missing = [w for w in extra.split() if w.lower() not in q.lower()]
            if missing:
                q = q + " " + " ".join(missing)
    except Exception:
        pass
    return q


def _fit_budget(parts: list[str], max_chars: int) -> str:
    """Distribuisce il budget di contesto tra le fonti invece di tagliare in
    coda: prima ogni fonte riceve una quota equa, così nessuna cartella viene
    scartata in silenzio solo perché arriva dopo la KB."""
    if not parts:
        return ""
    joined = "\n\n".join(parts)
    if len(joined) <= max_chars:
        return joined
    quota = max(700, max_chars // len(parts))
    trimmed = [p[:quota] for p in parts]
    return "\n\n".join(trimmed)[:max_chars]


def retrieve(dept: str, query: str, use_kb: bool = True, use_folder: bool = True,
             max_chars: int = 6000) -> str:
    """Contesto combinato secondo i toggle dell'utente: conoscenza ChromaDB del
    dipartimento e/o cartelle del dipartimento. Resiliente: in caso di errore in
    una sorgente, restituisce ciò che ha (la chat non deve fallire per il retrieval)."""
    parts = []
    if use_kb:
        kb_text, _ = kb_search(dept, query, n=5)
        if kb_text.strip():
            parts.append("[Conoscenza dipartimento — " + dept + "]\n" + kb_text)
    if use_folder:
        for p in store.department_folders(dept):
            # rerank semantico: riordina i candidati BM25 per pertinenza reale
            # (fallback automatico all'ordine BM25 se il modello non è disponibile)
            f_text, _ = folder_search(p, query, n=4, rerank=True)
            if f_text.strip():
                parts.append("[Cartella: " + p + "]\n" + f_text)
    return _fit_budget(parts, max_chars)


# ── Archivio ORIGINALI della Conoscenza ─────────────────────
# All'upload il file originale viene conservato nel volume dati (per
# dipartimento): le FONTI in chat diventano link scaricabili. Perimetro:
# serviti SOLO agli utenti del dipartimento, download tracciato in audit.
KB_FILES_DIR = Path(os.environ.get("APP_DATA_DIR", "/data")) / "kb_files"


def kb_file_path(dept: str, filename: str) -> Path:
    """Percorso dell'originale, a prova di traversal: directory per
    collezione di dipartimento + SOLO il nome base del file."""
    d = KB_FILES_DIR / store.collection_for_department(dept)
    d.mkdir(parents=True, exist_ok=True)
    name = Path(filename or "").name.strip()
    if not name or name in (".", ".."):
        name = "documento"
    return d / name


def kb_save_file(dept: str, filename: str, raw: bytes) -> None:
    try:
        kb_file_path(dept, filename).write_bytes(raw)
    except Exception as e:
        import sys
        print(f"[kb-file] impossibile archiviare {filename}: {e}", file=sys.stderr)


def kb_delete_file(dept: str, filename: str) -> None:
    try:
        p = kb_file_path(dept, filename)
        if p.is_file():
            p.unlink()
    except Exception:
        pass


def kb_links(dept: str, names: list) -> list:
    """Voci FONTI per i documenti KB: link all'originale se archiviato,
    solo il nome altrimenti (documenti caricati prima di questa funzione)."""
    from urllib.parse import quote
    out = []
    for n in names:
        name = n[0] if isinstance(n, (tuple, list)) else str(n)
        item = {"name": name, "kind": "kb"}
        if kb_file_path(dept, name).is_file():
            item["url"] = "/kb/file/" + quote(name)
        out.append(item)
    return out
