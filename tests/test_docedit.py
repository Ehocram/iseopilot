"""
Test EDITING IN PLACE dei documenti allegati (Incremento 10): rilevamento del
ramo, revisioni tracciate Word, applicazione PowerPoint/Excel, guardie
non-distruttive, flusso completo in chat.
Esecuzione:
    PYTHONPATH=. APP_DATA_DIR=./data_test python -m pytest tests/ -q
"""
import io
import os
import re
import zipfile

os.environ.setdefault("APP_DATA_DIR", "./data_test")
from cryptography.fernet import Fernet
os.environ.setdefault("APP_SECRET_KEY", Fernet.generate_key().decode())

import pytest
from fastapi.testclient import TestClient
from app.main import app
from app import store, auth, docedit

store.init_db()


def _mk_user(name, dept="IT"):
    store.create_user(name, auth.hash_password("Password123"), dept, is_admin=False)
    c = TestClient(app)
    c.post("/login", data={"username": name, "password": "Password123"},
           follow_redirects=False)
    return c


def _docx_bytes() -> bytes:
    from docx import Document
    d = Document()
    d.add_paragraph("Il prezzo indicato è di ")
    p = d.paragraphs[0]
    p.add_run("12.000 euro").bold = True
    p.add_run(" IVA esclusa.")
    d.add_paragraph("La garanzia è di 12 mesi dalla consegna.")
    t = d.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text = "Referente"
    t.rows[0].cells[1].text = "Mario Rossi"
    bio = io.BytesIO()
    d.save(bio)
    return bio.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Risultati 2025"
    s.shapes.placeholders[1].text_frame.text = "Fatturato in crescita del 12%"
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


def _xlsx_bytes(con_grafico=False) -> bytes:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Listino"
    ws["A1"], ws["B1"] = "Prodotto", "Prezzo"
    ws["A2"], ws["B2"] = "Cilindro X100", 120
    ws["B3"] = "=SUM(B2:B2)"
    if con_grafico:
        from openpyxl.chart import BarChart, Reference
        ch = BarChart()
        ch.add_data(Reference(ws, min_col=2, min_row=1, max_row=2))
        ws.add_chart(ch, "D2")
    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


def _att(name, ok=True):
    return {"id": "a" * 32, "name": name, "ok": ok}


# ── Rilevamento del ramo editing ────────────────────────────
def test_detect_edit_positivi():
    a = [_att("Offerta.docx")]
    casi = [
        "modifica il documento: il prezzo diventa 14.500 euro",
        "correggi l'allegato sostituendo il referente",
        "aggiorna il file con la garanzia a 24 mesi",
        "modificalo e restituiscimelo",
        "rivedi il documento allegato e rimandamelo",
        "please update the document: change the warranty to 24 months",
    ]
    for msg in casi:
        assert docedit.detect_edit_request(msg, a) is not None, msg


def test_detect_edit_negativi():
    a = [_att("Offerta.docx")]
    casi = [
        "creami un word con la sintesi di questo documento",   # creazione → docgen
        "genera una presentazione partendo dall'allegato",     # creazione → docgen
        "cosa dice il documento sulla garanzia?",              # domanda
        "riassumi l'allegato",                                 # domanda
        "che tempo fa domani?",                                # fuori tema
    ]
    for msg in casi:
        assert docedit.detect_edit_request(msg, a) is None, msg


def test_detect_edit_richiede_allegato_modificabile():
    msg = "modifica il documento e restituiscimelo"
    assert docedit.detect_edit_request(msg, []) is None
    assert docedit.detect_edit_request(msg, [_att("nota.pdf")]) is None   # pdf non editabile
    assert docedit.detect_edit_request(msg, [_att("nota.txt")]) is None
    assert docedit.detect_edit_request(msg, [_att("x.docx", ok=False)]) is None
    assert docedit.detect_edit_request(msg, [_att("ok.xlsx")]) is not None


def test_detect_edit_sceglie_allegato_nominato():
    atts = [_att("Listino.xlsx"), _att("Offerta.docx")]
    scelto = docedit.detect_edit_request("modifica il Listino: prezzo a 135", atts)
    assert scelto["name"] == "Listino.xlsx"
    # senza nome esplicito: l'ultimo caricato
    scelto = docedit.detect_edit_request("modifica il documento e rimandamelo", atts)
    assert scelto["name"] == "Offerta.docx"


# ── DOCX: revisioni tracciate ───────────────────────────────
def test_docx_revisioni_tracciate(tmp_path):
    src = tmp_path / "off.docx"
    src.write_bytes(_docx_bytes())
    plan = {"modifiche": [
        {"trova": "12.000 euro", "sostituisci": "14.500 euro", "motivo": "listino"},
        {"trova": "garanzia è di 12 mesi", "sostituisci": "garanzia è di 24 mesi", "motivo": "estensione"},
        {"trova": "Mario Rossi", "sostituisci": "Laura Bianchi", "motivo": "referente"},
    ]}
    out, ok, ko, avvisi = docedit.apply_docx(str(src), plan)
    assert len(ok) == 3 and not ko
    xml = zipfile.ZipFile(out).read("word/document.xml").decode()
    assert xml.count("<w:ins ") == 3 and xml.count("<w:del ") == 3
    assert 'w:author="ISEOPilot"' in xml
    # dentro w:del il testo DEVE stare in delText (contratto OOXML)
    for blocco in re.findall(r"<w:del .*?</w:del>", xml, re.S):
        assert "<w:delText" in blocco and "<w:t>" not in blocco
    # nuovo testo negli inserimenti, vecchio nelle cancellazioni
    ins = " ".join(re.findall(r"<w:ins .*?</w:ins>", xml, re.S))
    assert "14.500 euro" in ins and "Laura Bianchi" in ins
    dele = " ".join(re.findall(r"<w:del .*?</w:del>", xml, re.S))
    assert "12.000 euro" in dele and "Mario Rossi" in dele
    # formattazione ereditata: il prezzo era in grassetto
    blocco_prezzo = next(b for b in re.findall(r"<w:ins .*?</w:ins>", xml, re.S)
                         if "14.500" in b)
    assert "<w:b/>" in blocco_prezzo or "<w:b " in blocco_prezzo


def test_docx_ancoraggio_mancante_dichiarato(tmp_path):
    src = tmp_path / "off.docx"
    src.write_bytes(_docx_bytes())
    plan = {"modifiche": [
        {"trova": "clausola che non esiste", "sostituisci": "x", "motivo": "prova"},
        {"trova": "12.000 euro", "sostituisci": "13.000 euro", "motivo": "ok"},
    ]}
    out, ok, ko, _ = docedit.apply_docx(str(src), plan)
    assert len(ok) == 1 and len(ko) == 1
    assert "non trovato" in ko[0]["perche"]          # dichiarato, mai inventato


def test_docx_cancellazione_pura(tmp_path):
    src = tmp_path / "off.docx"
    src.write_bytes(_docx_bytes())
    plan = {"modifiche": [{"trova": " IVA esclusa.", "sostituisci": "", "motivo": "rimozione"}]}
    out, ok, ko, _ = docedit.apply_docx(str(src), plan)
    assert len(ok) == 1
    xml = zipfile.ZipFile(out).read("word/document.xml").decode()
    assert "<w:del " in xml and "<w:ins " not in xml   # solo cancellazione


def test_docx_ancoraggio_tollerante_a_spazi_e_apici(tmp_path):
    from docx import Document
    d = Document()
    d.add_paragraph("Il  contratto\u00a0prevede l\u2019adeguamento annuale.")
    src = tmp_path / "t.docx"
    d.save(str(src))
    plan = {"modifiche": [
        {"trova": "Il contratto prevede l'adeguamento annuale.",
         "sostituisci": "Il contratto prevede l'adeguamento biennale.", "motivo": "x"}]}
    out, ok, ko, _ = docedit.apply_docx(str(src), plan)
    assert len(ok) == 1 and not ko


# ── PPTX ────────────────────────────────────────────────────
def test_pptx_applica_e_registra_nelle_note(tmp_path):
    from pptx import Presentation
    src = tmp_path / "d.pptx"
    src.write_bytes(_pptx_bytes())
    plan = {"modifiche": [
        {"trova": "Risultati 2025", "sostituisci": "Risultati 2026", "motivo": "anno"},
        {"trova": "crescita del 12%", "sostituisci": "crescita del 18%", "motivo": "dato"},
        {"trova": "assente", "sostituisci": "x", "motivo": "prova"},
    ]}
    out, ok, ko, avvisi = docedit.apply_pptx(str(src), plan)
    assert len(ok) == 2 and len(ko) == 1
    prs = Presentation(out)
    sl = prs.slides[0]
    assert sl.shapes.title.text == "Risultati 2026"
    assert "18%" in sl.shapes.placeholders[1].text_frame.text
    note = sl.notes_slide.notes_text_frame.text
    assert "ISEOPilot" in note and "Risultati 2026" in note
    assert avvisi and "revisioni tracciate" in avvisi[0]


# ── XLSX ────────────────────────────────────────────────────
def test_xlsx_valore_e_commento(tmp_path):
    import openpyxl
    src = tmp_path / "l.xlsx"
    src.write_bytes(_xlsx_bytes())
    plan = {"modifiche": [
        {"foglio": "Listino", "cella": "B2", "sostituisci": "135", "motivo": "listino"},
        {"trova": "Cilindro X100", "sostituisci": "Cilindro X150", "motivo": "codice"},
    ]}
    out, ok, ko, avvisi = docedit.apply_xlsx(str(src), plan)
    assert len(ok) == 2 and not ko
    ws = openpyxl.load_workbook(out)["Listino"]
    assert ws["B2"].value == 135 and isinstance(ws["B2"].value, int)   # numero, non testo
    assert "120" in ws["B2"].comment.text                              # valore precedente
    assert ws["A2"].value == "Cilindro X150"


def test_xlsx_formule_mai_sovrascritte(tmp_path):
    import openpyxl
    src = tmp_path / "l.xlsx"
    src.write_bytes(_xlsx_bytes())
    plan = {"modifiche": [
        {"trova": "=SUM(B2:B2)", "sostituisci": "999", "motivo": "x"},
        {"foglio": "Listino", "cella": "B3", "sostituisci": "999", "motivo": "y"},
    ]}
    out, ok, ko, _ = docedit.apply_xlsx(str(src), plan)
    assert not ok and len(ko) == 2
    assert all("formula" in k["perche"] for k in ko)
    assert openpyxl.load_workbook(out)["Listino"]["B3"].value == "=SUM(B2:B2)"


def test_xlsx_con_grafici_rifiutato(tmp_path):
    src = tmp_path / "g.xlsx"
    src.write_bytes(_xlsx_bytes(con_grafico=True))
    assert docedit.xlsx_ha_grafici_o_immagini(str(src)) is True
    with pytest.raises(ValueError) as e:
        docedit.apply_plan("xlsx", str(src), {"modifiche": []})
    assert "grafici" in str(e.value) and "NON è stato toccato" in str(e.value)


# ── Piano: parsing e tetto ──────────────────────────────────
def test_build_plan_tetto_e_json_sporco(monkeypatch):
    from app import orchestrator
    tante = ",".join('{"trova": "a%d", "sostituisci": "b"}' % i for i in range(60))
    monkeypatch.setattr(docedit.orchestrator, "complete",
                        lambda s, u, st, max_tokens=4000, timeout=120:
                        '```json\n{"modifiche": [' + tante + '], "note": ""}\n```')
    plan = docedit.build_plan("docx", "richiesta", "testo", {"claude_api_key": "k"})
    assert len(plan["modifiche"]) == docedit.MAX_MODIFICHE
    assert plan["troncato"] is True


def test_build_plan_json_invalido_fail_loud(monkeypatch):
    monkeypatch.setattr(docedit.orchestrator, "complete",
                        lambda s, u, st, max_tokens=4000, timeout=120: "non è json")
    with pytest.raises(ValueError):
        docedit.build_plan("docx", "r", "t", {"claude_api_key": "k"})


def test_riepilogo_dichiara_applicate_e_non():
    testo = docedit.riepilogo(
        "Offerta (rev).docx", "docx",
        [{"trova": "12.000", "sostituisci": "14.500", "motivo": "listino"}],
        [{"trova": "clausola X", "perche": "testo non trovato nel documento"}],
        [], troncato=False, note="")
    assert "revisioni tracciate" in testo
    assert "Modifiche applicate (1)" in testo
    assert "NON applicate (1)" in testo and "clausola X" in testo
    assert "originale che hai caricato non è stato toccato" in testo


# ── Flusso completo in chat ─────────────────────────────────
def test_flusso_chat_modifica_documento_allegato(monkeypatch):
    c = _mk_user("edit@test")
    r = c.post("/api/attach", files=[("files", ("Offerta.docx", _docx_bytes(),
               "application/vnd.openxmlformats-officedocument.wordprocessingml.document"))])
    a = r.json()["attachments"][0]
    assert a["ok"] and a.get("editable") is True      # sorgente conservato

    monkeypatch.setattr(docedit.orchestrator, "complete",
                        lambda s, u, st, max_tokens=4000, timeout=120:
                        '{"modifiche": [{"trova": "12.000 euro", '
                        '"sostituisci": "14.500 euro", "motivo": "listino 2026"}], "note": ""}')
    r = c.post("/api/chat", json={"engine": "claude", "free_mode": True,
               "attachments": [a],
               "messages": [{"role": "user", "content":
                             "modifica il documento: il prezzo diventa 14.500 euro"}]})
    assert r.status_code == 200
    assert "revisioni tracciate" in r.text
    assert "Offerta (rev).docx" in r.text
    assert '"kind": "download"' in r.text
    # audit del ramo corretto
    azioni = {x["action"] for x in store.audit_query(username="edit@test")}
    assert "modifica_documento" in azioni


def test_chat_creazione_non_viene_dirottata_su_editing(monkeypatch):
    """Non-regressione: con un allegato modificabile, «creami un word» resta
    generazione (docgen), non editing in place."""
    c = _mk_user("edit2@test")
    r = c.post("/api/attach", files=[("files", ("Base.docx", _docx_bytes(),
               "application/octet-stream"))])
    a = r.json()["attachments"][0]
    visto = {}

    def fake_generate(fmt, req, ctx, st, hist="", templates=None, note_utente=""):
        visto["fmt"] = fmt
        import tempfile
        p = os.path.join(tempfile.gettempdir(), "g.docx")
        open(p, "wb").write(_docx_bytes())
        return p, "nuovo.docx"

    from app import main as main_mod
    monkeypatch.setattr(main_mod.docgen, "generate", fake_generate)
    r = c.post("/api/chat", json={"engine": "claude", "free_mode": True,
               "attachments": [a],
               "messages": [{"role": "user", "content":
                             "creami un word di sintesi da questo documento"}]})
    assert visto.get("fmt") == "docx"                 # ramo generazione
    assert "revisioni tracciate" not in r.text


def test_sorgente_scaduto_non_finge(monkeypatch):
    """Se i byte originali non sono più nel deposito, niente editing silenzioso."""
    c = _mk_user("edit3@test")
    finto = {"id": "b" * 32, "name": "Sparito.docx", "ok": True}
    r = c.post("/api/chat", json={"engine": "claude", "free_mode": True,
               "attachments": [finto],
               "messages": [{"role": "user", "content":
                             "modifica il documento e restituiscimelo"}]})
    assert r.status_code == 200
    assert "revisioni tracciate" not in r.text        # nessuna finzione di modifica


if __name__ == "__main__":
    import inspect
    import sys
    import tempfile
    from pathlib import Path

    class _MP:
        def setattr(self, obj, name, value):
            setattr(obj, name, value)

    fns = [f for n, f in sorted(globals().items())
           if n.startswith("test_") and inspect.isfunction(f)]
    failed = 0
    for f in fns:
        try:
            kwargs = {}
            sig = inspect.signature(f)
            td = None
            if "tmp_path" in sig.parameters:
                td = tempfile.TemporaryDirectory()
                kwargs["tmp_path"] = Path(td.name)
            if "monkeypatch" in sig.parameters:
                kwargs["monkeypatch"] = _MP()
            f(**kwargs)
            print(f"  PASS  {f.__name__}")
            if td:
                td.cleanup()
        except Exception as e:
            failed += 1
            print(f"  FAIL  {f.__name__}: {e}")
    print(f"\n{len(fns)-failed}/{len(fns)} test superati.")
    sys.exit(1 if failed else 0)
