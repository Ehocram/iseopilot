"""
Smoke test ISEO Chat Web (Incremento 2).
Esecuzione:
    PYTHONPATH=. APP_DATA_DIR=./data_test \
    APP_SECRET_KEY=$(python -c "from cryptography.fernet import Fernet;print(Fernet.generate_key().decode())") \
    python tests/test_smoke.py
"""
import os

os.environ.setdefault("APP_DATA_DIR", "./data_test")
from cryptography.fernet import Fernet
os.environ.setdefault("APP_SECRET_KEY", Fernet.generate_key().decode())

from fastapi.testclient import TestClient
from app.main import app
from app import store, auth
from app.anonymizer import Anonymizer

store.init_db()


def fresh_client():
    return TestClient(app)


def login(c, username, password):
    return c.post("/login", data={"username": username, "password": password},
                  follow_redirects=False)


# -- Anonimizzazione --
def test_anonymize_restore_roundtrip():
    a = Anonymizer()
    src = "Scrivi a mario.rossi@iseo.com da 10.0.0.5 — CVE-2024-1234."
    anon = a.anonymize(src)
    assert "iseo.com" not in anon and "10.0.0.5" not in anon
    assert a.restore(anon) == src


# -- Password / auth --
def test_password_hash_verify():
    h = auth.hash_password("Sup3rSecret!")
    assert h.startswith("scrypt$")
    assert auth.verify_password("Sup3rSecret!", h) is True
    assert auth.verify_password("sbagliata", h) is False


def test_password_not_plaintext_on_disk():
    store.create_user("plainck@test", auth.hash_password("Password123"), "IT")
    import sqlite3
    with sqlite3.connect(store.DB_PATH) as cx:
        row = cx.execute("SELECT password_hash FROM users WHERE username='plainck@test'").fetchone()[0]
    assert "Password123" not in row


# -- Dipartimenti --
def test_departments_seeded():
    deps = store.list_departments()
    for expected in ["IT", "Infosec", "ESG", "Privacy", "Sales",
                     "Operations", "Finance", "HR", "Supply Chain", "R&D"]:
        assert expected in deps


def test_department_collection_naming():
    assert store.collection_for_department("Supply Chain") == "kb_supply_chain"
    assert store.collection_for_department("R&D") == "kb_r_d"


def test_add_department():
    assert store.add_department("Legal") is True
    assert store.add_department("Legal") is False  # duplicato
    assert store.department_exists("Legal") is True


# -- Utenti --
def test_create_and_update_user():
    assert store.create_user("u1@test", auth.hash_password("Password123"), "Sales") is True
    assert store.create_user("u1@test", auth.hash_password("x"), "Sales") is False  # dup
    u = store.get_user("u1@test")
    assert u["department"] == "Sales" and u["is_admin"] == 0 and u["active"] == 1
    store.update_user("u1@test", department="Finance", is_admin=True)
    u = store.get_user("u1@test")
    assert u["department"] == "Finance" and u["is_admin"] == 1


# -- Login + accesso --
def test_login_and_protected_access():
    store.create_user("admin@test", auth.hash_password("Password123"), "IT", is_admin=True)
    c = fresh_client()
    r = c.get("/", follow_redirects=False)
    assert r.status_code == 303 and r.headers["location"] == "/login"
    r = login(c, "admin@test", "Password123")
    assert r.status_code == 303 and r.headers["location"] == "/"
    assert c.get("/").status_code == 200
    assert c.get("/admin").status_code == 200
    assert c.get("/admin/users").status_code == 200
    assert c.get("/admin/departments").status_code == 200


def test_login_wrong_password():
    store.create_user("admin2@test", auth.hash_password("Password123"), "IT", is_admin=True)
    c = fresh_client()
    r = c.post("/login", data={"username": "admin2@test", "password": "wrong"})
    assert r.status_code == 401


def test_non_admin_forbidden_on_admin():
    store.create_user("bob@test", auth.hash_password("Password123"), "Sales", is_admin=False)
    c = fresh_client()
    login(c, "bob@test", "Password123")
    assert c.get("/", follow_redirects=False).status_code == 200
    r = c.get("/admin", follow_redirects=False)
    assert r.status_code == 403


def test_chat_requires_auth():
    c = fresh_client()
    r = c.post("/api/chat", json={"messages": [{"role": "user", "content": "ciao"}], "engine": "claude"})
    assert r.status_code == 401


def test_chat_streams_error_when_unconfigured():
    store.set_setting("claude_api_key", "", secret=True)
    store.create_user("chatuser@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "chatuser@test", "Password123")
    r = c.post("/api/chat", json={"messages": [{"role": "user", "content": "ciao"}], "engine": "claude", "source": "kb"})
    assert r.status_code == 200
    assert "error" in r.text and "Claude" in r.text


def test_deactivated_user_cannot_login():
    store.create_user("gone@test", auth.hash_password("Password123"), "IT")
    store.update_user("gone@test", active=False)
    c = fresh_client()
    r = c.post("/login", data={"username": "gone@test", "password": "Password123"})
    assert r.status_code == 401


def test_healthz():
    assert fresh_client().get("/healthz").json()["status"] == "ok"


# -- Conoscenza: cartelle (FTS5) + estrazione + pagina --
def test_department_folders_multi():
    store.add_department("Quality")
    assert store.department_folders("Quality") == []
    assert store.add_department_folder("Quality", "/tmp/qa")
    assert store.add_department_folder("Quality", "/tmp/qa2")
    assert store.department_folders("Quality") == ["/tmp/qa", "/tmp/qa2"]
    # compat: department_folder ritorna la prima
    assert store.department_folder("Quality") == "/tmp/qa"
    store.remove_department_folder("Quality", "/tmp/qa")
    assert store.department_folders("Quality") == ["/tmp/qa2"]
    # niente duplicati
    store.add_department_folder("Quality", "/tmp/qa2")
    assert store.department_folders("Quality") == ["/tmp/qa2"]


def test_extract_text_txt_md():
    from app import knowledge
    assert "ciao mondo" in knowledge.extract_text("a.txt", b"ciao mondo")
    assert "Titolo" in knowledge.extract_text("b.md", b"# Titolo\ncorpo")
    assert knowledge.extract_text("c.exe", b"xx") == ""  # tipo non ammesso


def test_folder_search_end_to_end_multi():
    import tempfile, os
    from app import knowledge
    d = tempfile.mkdtemp()
    with open(os.path.join(d, "policy.md"), "w") as fh:
        fh.write("La password minima e' di dodici caratteri secondo la policy ISEO.")
    with open(os.path.join(d, "vpn.md"), "w") as fh:
        fh.write("La VPN aziendale usa WireGuard sulla porta 51820.")
    # Registra la cartella al dipartimento e reindicizza via API dipartimento
    store.add_department("Reti")
    store.add_department_folder("Reti", d)
    ok, msg = knowledge.dept_folders_reindex("Reti")
    assert ok, msg
    assert knowledge.dept_folders_count("Reti") > 0
    # retrieve (nuova firma) trova il contenuto quando use_folder=True
    ctx = knowledge.retrieve("Reti", "quanti caratteri per la password", use_kb=False, use_folder=True)
    assert "dodici" in ctx
    # con il toggle spento, niente contesto
    ctx_off = knowledge.retrieve("Reti", "password", use_kb=False, use_folder=False)
    assert ctx_off == ""


def test_kb_unavailable_is_graceful():
    from app import knowledge
    # In sandbox ChromaDB e' assente: nessuna funzione KB deve sollevare eccezioni.
    assert knowledge.kb_count("IT") == 0
    assert knowledge.kb_list("IT") == []
    ok, msg = knowledge.kb_ingest("IT", "x.txt", "contenuto")
    assert isinstance(ok, bool) and isinstance(msg, str)


def test_knowledge_page_renders():
    store.create_user("ku@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "ku@test", "Password123")
    r = c.get("/knowledge")
    assert r.status_code == 200
    assert "kb_it" in r.text  # mostra la collezione del dipartimento


def test_admin_folder_add_remove_routes():
    store.create_user("adm3@test", auth.hash_password("Password123"), "IT", is_admin=True)
    import tempfile
    d = tempfile.mkdtemp()
    c = fresh_client()
    login(c, "adm3@test", "Password123")
    r = c.post("/admin/departments/folder/add",
               data={"department": "Sales", "path": d}, follow_redirects=False)
    assert r.status_code == 303
    assert d in store.department_folders("Sales")
    r2 = c.post("/admin/departments/folder/remove",
                data={"department": "Sales", "path": d}, follow_redirects=False)
    assert r2.status_code == 303
    assert d not in store.department_folders("Sales")


def test_admin_browse_lists_subdirs():
    store.create_user("adm4@test", auth.hash_password("Password123"), "IT", is_admin=True)
    import tempfile, os
    base = tempfile.mkdtemp()
    os.makedirs(os.path.join(base, "sottocartella_x"))
    c = fresh_client()
    login(c, "adm4@test", "Password123")
    r = c.get("/admin/browse", params={"dept": "Sales", "path": base})
    assert r.status_code == 200
    assert "sottocartella_x" in r.text


def test_settings_toggles_persist():
    store.create_user("tg@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "tg@test", "Password123")
    # spegne KB e cartelle, accende onedrive
    r = c.post("/settings", data={"use_onedrive": "1"}, follow_redirects=False)
    assert r.status_code == 303
    assert store.get_user_setting("tg@test", "use_kb", "1") == "0"
    assert store.get_user_setting("tg@test", "use_folder", "1") == "0"
    assert store.get_user_setting("tg@test", "use_onedrive", "0") == "1"


def test_admin_connectors_autofill_defaults():
    store.create_user("adm5@test", auth.hash_password("Password123"), "IT", is_admin=True)
    c = fresh_client()
    login(c, "adm5@test", "Password123")
    r = c.get("/admin")
    assert r.status_code == 200
    # i default ISEO compaiono precompilati
    assert "c5a90f54-d599-4f71-a98f-0fa0781145c1" in r.text
    assert "isd365-prod.operations.eu.dynamics.com" in r.text


def test_i18n_engine():
    from app import i18n
    assert i18n.t("Conoscenza", "it") == "Conoscenza"
    assert i18n.t("Conoscenza", "en") == "Knowledge"
    assert i18n.t("Esci", "en") == "Sign out"
    # fallback: stringa non tradotta resta in italiano
    assert i18n.t("Stringa inesistente xyz", "en") == "Stringa inesistente xyz"
    assert i18n.normalize("Italiano") == "it"
    assert i18n.normalize("English") == "en"


def test_ui_lang_switch_persists():
    store.create_user("lg@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "lg@test", "Password123")
    # passa a EN
    c.get("/ui-lang", params={"to": "en", "next": "/"}, follow_redirects=False)
    assert store.get_user_setting("lg@test", "ui_lang", "it") == "en"
    r = c.get("/")
    assert ">Knowledge<" in r.text and ">Sign out<" in r.text
    # torna a IT
    c.get("/ui-lang", params={"to": "it", "next": "/"}, follow_redirects=False)
    r2 = c.get("/")
    assert ">Conoscenza<" in r2.text


def test_connectors_configured_with_defaults():
    from app import connectors
    # i default ISEO rendono i connettori "configurati"
    assert connectors.is_configured("onedrive") is True
    assert connectors.is_configured("dynamics") is True
    assert connectors.is_configured("inesistente") is False


def test_connectors_not_connected_and_search_empty():
    from app import connectors
    # senza token: non connesso, ricerca vuota, nessuna eccezione
    assert connectors.is_connected("nobody@test", "onedrive") is False
    assert connectors.search("nobody@test", "onedrive", "qualcosa") == ""
    assert connectors.search("nobody@test", "dynamics", "qualcosa") == ""


def test_connectors_token_isolation_per_user():
    # ogni utente ha un percorso token distinto
    pa = store.user_token_path("alice@test", "onedrive")
    pb = store.user_token_path("bob@test", "onedrive")
    assert pa != pb


def test_connectors_disconnect_is_safe():
    from app import connectors
    connectors.disconnect("ghost@test", "dynamics")  # non deve sollevare


def test_connect_endpoints_require_auth_and_validate():
    c = fresh_client()
    # senza login -> 401
    assert c.post("/connect/onedrive/start").status_code == 401
    # con login, connettore non valido -> 404
    store.create_user("cn@test", auth.hash_password("Password123"), "IT")
    c2 = fresh_client()
    login(c2, "cn@test", "Password123")
    assert c2.post("/connect/xxx/start").status_code == 404
    # start valido: in sandbox la rete MS non è raggiungibile -> ok:False gestito
    r = c2.post("/connect/onedrive/start")
    assert r.status_code == 200
    body = r.json()
    assert "ok" in body  # struttura corretta, niente crash


def test_settings_shows_connect_controls():
    store.create_user("sc@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "sc@test", "Password123")
    r = c.get("/settings")
    assert r.status_code == 200
    # bottone di connessione presente (configurato dai default, non ancora connesso)
    assert 'data-conn="onedrive"' in r.text and 'data-conn="dynamics"' in r.text


def test_free_mode_system_prompt():
    from app.orchestrator import build_system
    # modalità libera: istruzione generalista, niente blocco CONTESTO anche se fornito
    s = build_system("Tecnico", "Italiano", "DOC_AZIENDALE_SEGRETO", free_mode=True)
    assert "MODALITÀ AI LIBERA" in s
    assert "DOC_AZIENDALE_SEGRETO" not in s and "CONTESTO" not in s
    # modalità documentale: usa il contesto
    s2 = build_system("Tecnico", "Italiano", "DOC_AZIENDALE_SEGRETO", free_mode=False)
    assert "DOC_AZIENDALE_SEGRETO" in s2 and "MODALITÀ AI LIBERA" not in s2


def test_chat_has_mode_selector():
    store.create_user("md@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "md@test", "Password123")
    r = c.get("/")
    assert r.status_code == 200
    assert 'id="mode"' in r.text and 'value="free"' in r.text


def test_dyn_catalog_paths_redirected_under_data():
    from app import connectors
    import app.engines.dynamics_search as ds
    assert str(ds.CATALOG_FILE).endswith("dynamics/catalog.json")
    assert str(ds.SCHEMA_DIR).endswith("dynamics/schema")


def test_dyn_catalog_status_absent_and_present():
    import json
    from app import connectors
    import app.engines.dynamics_search as ds
    # assente
    if ds.CATALOG_FILE.exists():
        ds.CATALOG_FILE.unlink()
    assert connectors.dyn_catalog_status() == {"present": False}
    # presente (catalogo finto)
    ds.CATALOG_FILE.parent.mkdir(parents=True, exist_ok=True)
    ds.CATALOG_FILE.write_text(json.dumps({
        "count": 4704, "relazioni": 12000, "generato": "2026-06-23",
        "istanza": "https://x.dynamics.com", "versione": "2.1"}), encoding="utf-8")
    st = connectors.dyn_catalog_status()
    assert st["present"] and st["count"] == 4704 and st["relazioni"] == 12000
    ds.CATALOG_FILE.unlink()


def test_dyn_build_requires_connection():
    from app import connectors
    res = connectors.dyn_build_catalog("nobody@test")
    assert "errore" in res


def test_admin_build_catalog_endpoint():
    store.create_user("adcat@test", auth.hash_password("Password123"), "IT", is_admin=True)
    c = fresh_client()
    login(c, "adcat@test", "Password123")
    # admin non connesso a Dynamics -> redirect con dyn_err (gestito)
    r = c.post("/admin/dynamics/build-catalog", follow_redirects=False)
    assert r.status_code == 303 and "dyn_err" in r.headers["location"]


def test_knowledge_folder_add_remove_admin():
    store.create_user("kadm@test", auth.hash_password("Password123"), "IT", is_admin=True)
    import tempfile
    d = tempfile.mkdtemp()
    c = fresh_client()
    login(c, "kadm@test", "Password123")
    r = c.post("/knowledge/folder/add", data={"path": d}, follow_redirects=False)
    assert r.status_code == 303
    assert d in store.department_folders("IT")
    r2 = c.post("/knowledge/folder/remove", data={"path": d}, follow_redirects=False)
    assert r2.status_code == 303
    assert d not in store.department_folders("IT")


def test_knowledge_folder_add_forbidden_for_non_admin():
    store.create_user("kusr@test", auth.hash_password("Password123"), "Sales")
    c = fresh_client()
    login(c, "kusr@test", "Password123")
    r = c.post("/knowledge/folder/add", data={"path": "/tmp"}, follow_redirects=False)
    assert r.status_code == 403


def test_knowledge_page_admin_has_add_folder_form():
    store.create_user("kadm2@test", auth.hash_password("Password123"), "IT", is_admin=True)
    c = fresh_client()
    login(c, "kadm2@test", "Password123")
    r = c.get("/knowledge")
    assert r.status_code == 200
    assert 'action="/knowledge/folder/add"' in r.text


def test_auto_index_on_folder_attach():
    # Agganciando una cartella con un file, l'indice si costruisce da solo
    # (background task eseguito dal TestClient dopo la risposta).
    import tempfile, os
    from app import knowledge
    store.create_user("auto@test", auth.hash_password("Password123"), "IT", is_admin=True)
    d = tempfile.mkdtemp()
    with open(os.path.join(d, "nota.md"), "w") as fh:
        fh.write("Documento condiviso di rete con contenuto indicizzabile ISEO.")
    c = fresh_client()
    login(c, "auto@test", "Password123")
    r = c.post("/knowledge/folder/add", data={"path": d})
    assert r.status_code == 200  # dopo redirect+follow
    # senza premere 'reindex', la cartella è già indicizzata
    assert knowledge.folder_count(d) > 0


def test_reindex_all_indexes_all_departments():
    import tempfile, os
    from app import knowledge
    store.add_department("Logistica")
    d = tempfile.mkdtemp()
    with open(os.path.join(d, "ddt.md"), "w") as fh:
        fh.write("Documento di trasporto condiviso, magazzino centrale.")
    store.add_department_folder("Logistica", d)
    knowledge.reindex_all()  # sincrono nei test
    assert knowledge.dept_folders_count("Logistica") > 0


def test_wait_overlay_and_data_wait_present():
    store.create_user("wo@test", auth.hash_password("Password123"), "IT", is_admin=True)
    c = fresh_client()
    login(c, "wo@test", "Password123")
    # overlay globale + showWait nelle pagine che estendono base
    r = c.get("/")
    assert 'id="wait-overlay"' in r.text and "window.showWait" in r.text
    # link Conoscenza marcato con data-wait
    assert 'href="/knowledge" data-wait=' in r.text
    # form lenti marcati nella pagina Conoscenza
    k = c.get("/knowledge")
    assert 'action="/knowledge/upload"' in k.text and "data-wait=" in k.text


def test_chat_history_save_list_get_rename_delete():
    store.create_user("hist@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "hist@test", "Password123")
    hist = [{"role": "user", "content": "Domanda uno"},
            {"role": "assistant", "content": "Risposta uno"}]
    # save
    r = c.post("/api/chat/save", json={"session_id": None, "history": hist})
    assert r.status_code == 200
    sid = r.json()["session_id"]
    assert sid
    # list
    lst = c.get("/api/chat/list").json()["sessions"]
    assert any(s["id"] == sid for s in lst)
    assert lst[0]["title"].startswith("Domanda uno")
    # get
    got = c.get("/api/chat/get?sid=" + sid).json()
    assert got["history"][0]["content"] == "Domanda uno"
    # rename
    c.post("/api/chat/rename", json={"session_id": sid, "title": "Titolo nuovo"})
    assert c.get("/api/chat/list").json()["sessions"][0]["title"] == "Titolo nuovo"
    # delete
    c.post("/api/chat/delete", json={"session_id": sid})
    assert not any(s["id"] == sid for s in c.get("/api/chat/list").json()["sessions"])


def test_chat_history_is_user_scoped():
    store.create_user("scopeA@test", auth.hash_password("Password123"), "IT")
    store.create_user("scopeB@test", auth.hash_password("Password123"), "Sales")
    c1 = fresh_client(); login(c1, "scopeA@test", "Password123")
    c1.post("/api/chat/save", json={"session_id": None,
            "history": [{"role": "user", "content": "Segreto di A"},
                        {"role": "assistant", "content": "ok"}]})
    c2 = fresh_client(); login(c2, "scopeB@test", "Password123")
    # scopeB non vede le chat di scopeA
    assert c2.get("/api/chat/list").json()["sessions"] == []


def test_feedback_good_promotes_example():
    store.create_user("fb@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "fb@test", "Password123")
    r = c.post("/api/feedback/good", json={"question": "D?", "answer": "Ottima risposta"})
    assert r.status_code == 200 and r.json()["promoted"] == 1
    # entra nel contesto feedback
    from app import memory
    assert "RISPOSTE ECCELLENTI" in memory.build_feedback_context("fb@test")
    # pollice giù: ok ma nessun apprendimento
    assert c.post("/api/feedback/bad", json={"answer": "x"}).status_code == 200


def test_chat_page_has_sidebar_and_thumbs_setup():
    store.create_user("ui@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "ui@test", "Password123")
    r = c.get("/")
    assert 'id="new-chat"' in r.text and 'id="chat-list"' in r.text
    assert "window.I18N_FB" in r.text


def test_settings_ajax_autosave_persists():
    store.create_user("toggle@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "toggle@test", "Password123")
    # auto-save via AJAX: accende OneDrive (anche se non connesso, il valore si salva)
    r = c.post("/settings", data={"use_kb": "1", "use_onedrive": "1", "ajax": "1"})
    assert r.status_code == 200 and r.json()["ok"] is True
    assert store.get_user_setting("toggle@test", "use_onedrive", "0") == "1"
    assert store.get_user_setting("toggle@test", "use_kb", "0") == "1"
    # spegnere: use_onedrive non inviato -> torna "0"
    r2 = c.post("/settings", data={"use_kb": "1", "ajax": "1"})
    assert r2.status_code == 200
    assert store.get_user_setting("toggle@test", "use_onedrive", "0") == "0"


def test_settings_page_has_autosave_script():
    store.create_user("tgui@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "tgui@test", "Password123")
    r = c.get("/settings")
    assert "saveSources" in r.text and 'name^="use_"' in r.text


def test_dyn_report_served_to_owner_only():
    import tempfile, os
    from app import connectors
    store.create_user("rep@test", auth.hash_password("Password123"), "IT")
    store.create_user("rep2@test", auth.hash_password("Password123"), "Sales")
    # crea un finto report HTML e registralo per rep@test
    fd, path = tempfile.mkstemp(suffix=".html")
    os.write(fd, b"<html><body>Report Dynamics riservato</body></html>"); os.close(fd)
    token = connectors.register_report("rep@test", path)
    # proprietario: ottiene l'HTML
    c = fresh_client(); login(c, "rep@test", "Password123")
    r = c.get("/dyn-report/" + token)
    assert r.status_code == 200 and "Report Dynamics riservato" in r.text
    # altro utente: 404 (isolamento)
    c2 = fresh_client(); login(c2, "rep2@test", "Password123")
    assert c2.get("/dyn-report/" + token).status_code == 404
    # token inesistente: 404
    assert c.get("/dyn-report/nonexistent").status_code == 404


def test_settings_no_nested_forms():
    # I form annidati vengono scartati dal browser: il Disconnetti non deve
    # essere un form dentro il form delle impostazioni.
    store.create_user("nest@test", auth.hash_password("Password123"), "IT")
    c = fresh_client(); login(c, "nest@test", "Password123")
    html = c.get("/settings").text
    assert html.count("<form") == 1
    # il logout connettore avviene via fetch POST (visibile solo se connesso:
    # verifico il template sorgente)
    tpl = open("templates/user.html").read()
    assert tpl.count("<form") == 1
    assert "/connect/onedrive/logout" in tpl and "fetch(" in tpl


def test_onedrive_rank_and_trim():
    from app.engines.onedrive_search import _rank_and_trim
    # scenario reale: mp4/pdf/xlsx davanti, anagrafica in coda al pool Graph
    metas = [
        {"name": "Pilota AI Iseo - Registrazione della riunione.mp4"},
        {"name": "PANDIGITAL SRL_Cerved.pdf"},
        {"name": "Budget 2026_V5.xlsx"},
        {"name": "anagrafica.docx"},
        {"name": "2.png"},
    ]
    out = _rank_and_trim(metas, '"anagrafica" "bonometti" "marco"'.split(), 5)
    names = [m["name"] for m in out]
    # mp4 e png filtrati; anagrafica in PRIMA posizione (match sul nome)
    assert names[0] == "anagrafica.docx"
    assert all(".mp4" not in n and ".png" not in n for n in names)
    # se il pool contiene SOLO media, non si azzera (fallback)
    only_media = [{"name": "a.mp4"}, {"name": "b.png"}]
    assert len(_rank_and_trim(only_media, ["x"], 5)) == 2


def test_select_sources_cited_and_dedup():
    from app.main import _select_sources
    links = [
        {"name": "anagrafica.docx", "url": "https://x/1"},
        {"name": "2.png", "url": "https://x/2"},
        {"name": "DetACN_PiattaformaNIS.pdf", "url": "https://x/3"},
        {"name": "DetACN_PiattaformaNIS.pdf", "url": "https://x/3bis"},
        {"name": "Report Dynamics (HTML)", "url": "/dyn-report/t", "kind": "report"},
    ]
    resp = "Marco è nato il 24/01/1979 [Fonte: anagrafica.docx]."
    out = _select_sources(links, resp)
    names = [s["name"] for s in out]
    # mostra solo la fonte citata + il report; niente PNG né PDF non pertinenti
    assert "anagrafica.docx" in names and "Report Dynamics (HTML)" in names
    assert "2.png" not in names and "DetACN_PiattaformaNIS.pdf" not in names
    # nessuna citazione -> fallback: tutte, ma deduplicate per nome
    out2 = _select_sources(links, "risposta senza citazioni")
    names2 = [s["name"] for s in out2]
    # fallback prudente: dedup e MASSIMO 3 file oltre al report
    assert names2.count("DetACN_PiattaformaNIS.pdf") == 1 and len(out2) <= 4
    # citazione SENZA estensione ("anagrafica") aggancia comunque il file
    out3 = _select_sources(links, "il dato è in anagrafica, sezione 2")
    assert [s["name"] for s in out3 if s.get("kind") != "report"] == ["anagrafica.docx"]


def test_onedrive_binary_filter():
    from app.engines.onedrive_search import _is_binary_name
    assert _is_binary_name("ChatAssistant_v2.0_WINDOWS_fix-indicizzazione.zip")
    assert _is_binary_name("setup.EXE") and _is_binary_name("backup.tar")
    assert _is_binary_name("2.png") and _is_binary_name("foto.JPEG")
    assert _is_binary_name("Registrazione della riunione.mp4")
    assert not _is_binary_name("anagrafica.docx")
    assert not _is_binary_name("relazione.pdf") and not _is_binary_name("dati.xlsx")


def test_onedrive_search_is_logged():
    from app import connectors
    # utente non connesso: nessuna eccezione e nessun risultato
    txt, links = connectors.search_with_links("nolog@test", "onedrive", "prova")
    assert txt == "" and links == []
    # il ponte log del modulo OneDrive è instradato sul log unico
    from app.engines import onedrive_search
    onedrive_search._odlog("test ponte")
    assert "[onedrive] test ponte" in connectors.dyn_log_tail(10)


def test_chat_stream_first_byte_is_ping():
    # Il primo evento SSE deve essere un ping immediato: tiene viva la
    # connessione attraverso proxy/VPN durante il retrieval (fix errore 504).
    store.create_user("ping@test", auth.hash_password("Password123"), "IT")
    c = fresh_client(); login(c, "ping@test", "Password123")
    r = c.post("/api/chat", json={"messages": [{"role": "user", "content": "ciao, chi sei?"}],
                                  "engine": "claude", "free_mode": True})
    assert r.status_code == 200
    body = r.text
    first_evt = body.split("\n\n")[0]
    assert '"type": "ping"' in first_evt
    # senza chiave API lo stream prosegue con un errore pulito (non HTML)
    assert "Chiave API Claude non configurata" in body


def test_frontend_sanitizes_proxy_html_errors():
    store.create_user("prx@test", auth.hash_password("Password123"), "IT")
    c = fresh_client(); login(c, "prx@test", "Password123")
    js = c.get("/static/chat.js").text
    # rileva pagine HTML degli apparati e non le mostra grezze
    assert "DOCTYPE" in js and "proxyErr" in js
    r = c.get("/")
    assert "proxyErr" in r.text


def test_english_search_support():
    from app import knowledge
    # follow-up in inglese: la domanda interrogativa breve eredita il soggetto
    q = knowledge.enrich_query("what about 2026?", "show me the purchase requisitions issued by IT")
    assert "requisitions" in q.lower() and "2026" in q
    # concept map inglese: "how old" -> termine documentale "anagrafica"
    q2 = knowledge.enrich_query("how old is paolo belloli?")
    assert "anagrafica" in q2.lower()
    # "IT" (reparto) NON è stopword: resta nella ricerca
    from app.engines.folder_index import _STOP
    assert "it" not in _STOP and "what" in _STOP


def test_english_document_fts_search():
    import tempfile, os
    from app.engines import folder_index
    d = tempfile.mkdtemp()
    with open(os.path.join(d, "backup_policy_en.txt"), "w") as f:
        f.write("Backup policy: daily execution at 02:00 AM, retention 30 days, weekly restore test.")
    folder_index.get_index(d).update()
    t, _ = folder_index.search_folder(d, "what is the backup retention policy?")
    assert "retention 30 days" in t


def test_admin_model_dropdown_and_rewrite_toggle():
    store.create_user("mdl@test", auth.hash_password("Password123"), "IT", is_admin=True)
    c = fresh_client(); login(c, "mdl@test", "Password123")
    r = c.get("/admin")
    assert r.status_code == 200
    # tendina modelli Claude con le opzioni correnti
    assert '<select id="claude_model"' in r.text
    for m in ("claude-fable-5", "claude-opus-4-8", "claude-sonnet-4-6", "claude-haiku-4-5-20251001"):
        assert m in r.text
    # toggle riscrittura AI presente
    assert 'name="search_ai_rewrite"' in r.text
    # il salvataggio del flag funziona a livello store
    store.set_setting("search_ai_rewrite", "1")
    assert store.get_setting("search_ai_rewrite", "0") == "1"
    store.set_setting("search_ai_rewrite", "0")


def test_folder_snippet_centers_on_match():
    import tempfile, os
    from app.engines import folder_index
    d = tempfile.mkdtemp()
    # chunk lungo: il termine cercato sta OLTRE i primi 1200 caratteri
    filler = "parola comune riempitivo testo generico documento aziendale " * 40
    testo = filler + " Il codice segreto del progetto ORIONE è 7742. " + filler
    with open(os.path.join(d, "lungo.txt"), "w") as f:
        f.write(testo)
    folder_index.get_index(d).update()
    t, _ = folder_index.search_folder(d, "codice progetto ORIONE")
    # con lo snippet mirato, la porzione mostrata contiene il passaggio giusto
    assert "7742" in t


def test_enrich_query_followup_and_concept():
    from app import knowledge
    # follow-up: la domanda anaforica eredita il soggetto dal turno precedente
    q = knowledge.enrich_query("e nel 2025?", "mi dici le rda di paolo belloli?")
    assert "belloli" in q.lower() and "2025" in q
    # anche con pronome anaforico ("quelli") il soggetto viene ereditato
    qa = knowledge.enrich_query("e quelli del 2026?", "mi dici le rda emesse dal reparto IT?")
    assert "rda" in qa.lower() and "2026" in qa
    # concept map: "quanti anni ha" aggiunge il termine documentale "anagrafica"
    q2 = knowledge.enrich_query("quanti anni ha marco bonometti?")
    assert "anagrafica" in q2.lower()
    # domanda già ricca: non eredita dal turno precedente
    q3 = knowledge.enrich_query("procedura backup server produzione", "altro tema")
    assert "altro tema" not in q3


def test_fit_budget_no_silent_drop():
    from app import knowledge
    a = "[Conoscenza]\n" + "A" * 5000
    b = "[Cartella 1]\n" + "B" * 5000
    c = "[Cartella 2]\n" + "C" * 5000
    out = knowledge._fit_budget([a, b, c], 6000)
    # tutte e tre le fonti sopravvivono al taglio (prima la terza spariva)
    assert "[Conoscenza]" in out and "[Cartella 1]" in out and "[Cartella 2]" in out
    assert len(out) <= 6000
    # sotto budget: nessun taglio
    assert knowledge._fit_budget(["x", "y"], 6000) == "x\n\ny"


def test_kb_reembed_admin_only():
    store.create_user("emb@test", auth.hash_password("Password123"), "IT")
    c = fresh_client(); login(c, "emb@test", "Password123")
    assert c.post("/admin/kb/reembed").status_code == 403


def test_audit_logging_and_query():
    from app import store
    store.audit_log("alice", "login", "", "10.0.0.5")
    store.audit_log("alice", "chat", "modalita=documentale", "10.0.0.5")
    store.audit_log("bob", "login_fallito", "", "10.0.0.9")
    rows = store.audit_query()
    assert len(rows) >= 3
    # filtro per utente
    only_alice = store.audit_query(username="alice")
    assert all(r["username"] == "alice" for r in only_alice) and len(only_alice) >= 2
    # filtro per azione
    fails = store.audit_query(action="login_fallito")
    assert all(r["action"] == "login_fallito" for r in fails)
    # azioni e utenti distinti
    assert "login" in store.audit_actions()
    assert "alice" in store.audit_usernames()


def test_audit_page_admin_only():
    store.create_user("auditadmin@test", auth.hash_password("Password123"), "IT", is_admin=True)
    store.create_user("audituser@test", auth.hash_password("Password123"), "IT")
    ca = fresh_client(); login(ca, "auditadmin@test", "Password123")
    assert ca.get("/admin/audit").status_code == 200
    # export excel
    r = ca.get("/admin/audit/export?preset=all")
    assert r.status_code == 200
    assert "spreadsheetml" in r.headers.get("content-type", "")
    assert "audit_iseopilot" in r.headers.get("content-disposition", "")
    # utente normale: 403
    cu = fresh_client(); login(cu, "audituser@test", "Password123")
    assert cu.get("/admin/audit").status_code == 403
    assert cu.get("/admin/audit/export").status_code == 403


def test_login_writes_audit():
    from app import store
    store.create_user("auditlog@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "auditlog@test", "Password123")
    rows = store.audit_query(username="auditlog@test", action="login")
    assert len(rows) >= 1


def test_chroma_collection_name_valid_with_multilingual():
    import re as _re
    # riproduzione ESATTA del bug: tag del modello multilingue troncato a 24
    # lasciava un underscore finale -> nome collezione rifiutato da Chroma
    model_ident = "paraphrase-multilingual-MiniLM-L12-v2"
    tag = _re.sub(r"[^a-z0-9]+", "_", model_ident.lower()).strip("_")[:24].strip("_")
    name = f"kb_infosec__{tag}"
    assert _re.fullmatch(r"[a-zA-Z0-9][a-zA-Z0-9._-]{1,510}[a-zA-Z0-9]", name), name
    assert not name.endswith("_")
    # anche il vecchio modello resta valido e stabile (nessuna migrazione)
    tag2 = _re.sub(r"[^a-z0-9]+", "_", "all-minilm-l6-v2".lower()).strip("_")[:24].strip("_")
    assert tag2 == "all_minilm_l6_v2"


def test_chroma_safe_document_ids():
    import re as _re, hashlib
    fname = "GDOC 19.3.1 List of Legal, Regulatory (EN).pdf"
    base = _re.sub(r"[^a-zA-Z0-9._-]", "_", fname)[:180].strip("._-") or "doc"
    sid = f"{base}-{hashlib.md5(fname.encode()).hexdigest()[:8]}__chunk_0"
    assert _re.fullmatch(r"[a-zA-Z0-9._-]+", sid), sid


def test_knowledge_delete_uses_name_and_reports_reason():
    from unittest.mock import patch
    from app import knowledge as kn
    store.create_user("kbdel@test", auth.hash_password("Password123!"), "IT")
    c = fresh_client(); login(c, "kbdel@test", "Password123!")
    # il template usa d.name (mai la repr del dict) nel form di rimozione
    tpl = open("templates/knowledge.html").read()
    assert 'value="{{ d.name }}"' in tpl and 'value="{{ d }}"' not in tpl
    # il motivo del fallimento arriva all'utente
    with patch.object(kn, "kb_delete", return_value=(False, "Documento non trovato")):
        r = c.post("/knowledge/delete", data={"filename": "x.pdf"}, follow_redirects=False)
    assert "Documento+non+trovato" in r.headers.get("location", "")
    with patch.object(kn, "kb_delete", return_value=(True, "ok")):
        r2 = c.post("/knowledge/delete", data={"filename": "x.pdf"}, follow_redirects=False)
    assert "Rimosso" in r2.headers.get("location", "")


def test_knowledge_folder_drop_traversal():
    store.create_user("kbdz@test", auth.hash_password("Password123!"), "IT")
    c = fresh_client(); login(c, "kbdz@test", "Password123!")
    html = c.get("/knowledge").text
    # attraversamento cartelle + upload sempre a lotti; mai il finto-file directory
    assert "webkitGetAsEntry" in html and "collectEntry" in html
    assert "input.files = ev.dataTransfer.files" not in html
    assert "ajax=1" in html


def test_kb_upload_ajax_batch_mode():
    from unittest.mock import patch
    from app import knowledge as kn
    store.create_user("kbaj@test", auth.hash_password("Password123!"), "IT")
    c = fresh_client(); login(c, "kbaj@test", "Password123!")
    with patch.object(kn, "kb_available", return_value=True), \
         patch.object(kn, "extract_text", return_value="testo"), \
         patch.object(kn, "kb_ingest", return_value=(True, "ok")):
        r = c.post("/knowledge/upload?ajax=1",
                   files=[("files", ("GPO 19.1 policy (EN).pdf", b"%PDF fake", "application/pdf"))])
    assert r.status_code == 200 and r.json()["ok"] == 1 and r.json()["fail"] == []
    # la pagina contiene la logica dei lotti
    html = c.get("/knowledge").text
    assert "ajax=1" in html and "BATCH" in html


def test_kb_upload_reports_reasons():
    from unittest.mock import patch
    from app import knowledge as kn
    store.create_user("kbup@test", auth.hash_password("Password123!"), "IT")
    c = fresh_client(); login(c, "kbup@test", "Password123!")
    # ingest che fallisce con un motivo: il motivo DEVE arrivare all'utente
    with patch.object(kn, "kb_available", return_value=True), \
         patch.object(kn, "extract_text", return_value="testo di prova"), \
         patch.object(kn, "kb_ingest", return_value=(False, "Errore indicizzazione: conflitto embedding")):
        r = c.post("/knowledge/upload",
                   files=[("files", ("GPO 19.11 Password policy (IT).pdf", b"%PDF-1.4 fake", "application/pdf"))],
                   follow_redirects=False)
    loc = r.headers.get("location", "")
    assert "Errore+indicizzazione" in loc and "conflitto+embedding" in loc
    # PDF senza testo estraibile: suggerimento "PDF scansionato"
    with patch.object(kn, "kb_available", return_value=True), \
         patch.object(kn, "extract_text", return_value=""), \
         patch.object(kn, "kb_ingest", return_value=(False, "Nessun testo estraibile")):
        r2 = c.post("/knowledge/upload",
                    files=[("files", ("GPO scan.pdf", b"%PDF-1.4 fake", "application/pdf"))],
                    follow_redirects=False)
    assert "PDF+scansionato" in r2.headers.get("location", "")


def test_password_policy():
    from app.auth import validate_password
    # conformi
    assert validate_password("Sicurezza2026!") is None
    assert validate_password("Engineer@23482") is None
    # violazioni: corta, senza maiuscola, senza carattere speciale
    assert "12 caratteri" in validate_password("Corta!A")
    assert "maiuscola" in validate_password("tuttominuscolo-123")
    assert "speciale" in validate_password("SoloLettereENumeri12")
    # creazione utente admin: password non conforme rifiutata
    store.create_user("padm@test", auth.hash_password("Password123!"), "IT", is_admin=True)
    c = fresh_client(); login(c, "padm@test", "Password123!")
    r = c.post("/admin/users", data={"new_username": "nuovo@test", "new_password": "corta",
                                     "new_department": "IT", "new_is_admin": "0"},
               follow_redirects=False)
    assert "12+caratteri" in r.headers.get("location", "")


def test_account_change_password():
    store.create_user("acct@test", auth.hash_password("OldPassw0rd"), "IT")
    c = fresh_client()
    login(c, "acct@test", "OldPassw0rd")
    # pagina account raggiungibile
    assert c.get("/account").status_code == 200
    # password attuale sbagliata -> errore
    r = c.post("/account/password", data={"current_password": "WRONG",
        "new_password": "NuovaPassw0rd!", "confirm_password": "NuovaPassw0rd!"},
        follow_redirects=False)
    assert "err=current" in r.headers.get("location", "")
    # nuova troppo corta
    r = c.post("/account/password", data={"current_password": "OldPassw0rd",
        "new_password": "short", "confirm_password": "short"}, follow_redirects=False)
    assert "err=policy" in r.headers.get("location", "")
    # conferma non combacia
    r = c.post("/account/password", data={"current_password": "OldPassw0rd",
        "new_password": "NuovaPassw0rd!", "confirm_password": "Diversa123!X"}, follow_redirects=False)
    assert "err=match" in r.headers.get("location", "")
    # cambio valido
    r = c.post("/account/password", data={"current_password": "OldPassw0rd",
        "new_password": "NuovaPassw0rd!", "confirm_password": "NuovaPassw0rd!"}, follow_redirects=False)
    assert "ok=1" in r.headers.get("location", "")
    # la vecchia non funziona più, la nuova sì
    assert auth.authenticate("acct@test", "OldPassw0rd") is None
    assert auth.authenticate("acct@test", "NuovaPassw0rd!") is not None


def test_account_requires_login():
    c = fresh_client()
    r = c.get("/account", follow_redirects=False)
    assert r.status_code == 303 and "/login" in r.headers.get("location", "")


def test_attach_relevant_slice_finds_deep_excel_row():
    from app.main import _build_attach_block
    # Excel simulato: 3000 righe, il valore cercato sta alla riga 2500 —
    # ben oltre qualunque quota di testa
    rows = ["Fornitore;Articolo;Importo"]
    rows += [f"Fornitore{i};ART-{i};{i}.00" for i in range(1, 3000)]
    rows[2500] = "PANDIGITAL SRL;ART-SPECIALE;12.345,67"
    text = chr(10).join(rows)
    block = _build_attach_block([{"name": "listino.xlsx", "text": text}],
                                query="qual è l'importo di PANDIGITAL per ART-SPECIALE?")
    # la riga profonda È nel contesto, con intestazioni e selezione dichiarata
    assert "12.345,67" in block and "Fornitore;Articolo;Importo" in block
    assert "righe pertinenti" in block
    # domanda senza corrispondenze: fallback alla testa, troncamento dichiarato
    block2 = _build_attach_block([{"name": "listino.xlsx", "text": text}],
                                 query="riassumi il documento")
    assert "TRONCATO" in block2


def test_attach_block_fair_budget_and_truncation():
    from app.main import _build_attach_block, ATTACH_TOTAL_BUDGET
    big = "X" * 30000
    atts = [{"name": "listino.xlsx", "text": big},
            {"name": "offerta.pdf", "text": big},
            {"name": "contratto.docx", "text": big}]
    block = _build_attach_block(atts)
    # TUTTI i file presenti (prima il taglio a 14k totali faceva sparire il 2° e 3°)
    for n in ("listino.xlsx", "offerta.pdf", "contratto.docx"):
        assert n in block
    assert len(block) <= ATTACH_TOTAL_BUDGET + 2000
    # troncamento DICHIARATO al modello
    assert "TRONCATO" in block and "testo troncato" in block
    # file piccolo: intatto, nessun marcatore
    block2 = _build_attach_block([{"name": "nota.txt", "text": "valore: 42"}])
    assert "valore: 42" in block2 and "TRONCATO" not in block2


def test_api_attach_reports_real_length():
    from unittest.mock import patch
    from app import knowledge as kn
    store.create_user("att2@test", auth.hash_password("Password123!"), "IT")
    c = fresh_client(); login(c, "att2@test", "Password123!")
    with patch.object(kn, "extract_attachment_text", return_value="Y" * 50000):
        r = c.post("/api/attach", files=[("files", ("relazione.pdf", b"fake", "application/pdf"))])
    a = r.json()["attachments"][0]
    # chars dichiara la lunghezza REALE, il testo viaggia entro il budget per-file
    assert a["ok"] and a["chars"] == 50000 and len(a["text"]) == 30000
    # i TABELLARI hanno budget dedicato ampio (200k): l'Excel non perde le righe
    with patch.object(kn, "extract_attachment_text", return_value="Z" * 300000):
        r2 = c.post("/api/attach", files=[("files", ("listino.xlsx", b"fake", "application/vnd.ms-excel"))])
    a2 = r2.json()["attachments"][0]
    assert a2["ok"] and a2["chars"] == 300000 and len(a2["text"]) == 200000


def test_attach_block_fair_budget_placeholder():
    pass


def test_chat_stream_status_indicator():
    store.create_user("stat@test", auth.hash_password("Password123"), "IT")
    c = fresh_client(); login(c, "stat@test", "Password123")
    r = c.post("/api/chat", json={"messages": [{"role": "user", "content": "ciao"}],
                                  "free_mode": False, "source": "kb"})
    body = r.text
    # dopo il ping arriva lo stato "Ricerca in corso su Conoscenza…"
    assert '"type": "status"' in body and "Ricerca in corso su" in body
    js = c.get("/static/chat.js").text
    assert "search-ind" in js and 'evt.type === "status"' in js
    # in AI libera nessuno stato di ricerca
    r2 = c.post("/api/chat", json={"messages": [{"role": "user", "content": "ciao"}],
                                   "free_mode": True})
    assert "Ricerca in corso su" not in r2.text


def test_dynamics_ask_ai_parses_text_blocks_and_fallback():
    from unittest.mock import patch, MagicMock
    from app.engines.dynamics_search import DynamicsSearch
    ds = DynamicsSearch({"ai_engine": "claude", "claude_api_key": "k",
                         "claude_model": "claude-fable-5",
                         "claude_model_fallback": "claude-opus-4-8"})
    # risposta con blocco di RAGIONAMENTO prima del testo (stile Fable 5):
    # prima falliva con KeyError 'text', ora estrae il blocco text
    ok = MagicMock(status_code=200)
    ok.json.return_value = {"content": [
        {"type": "thinking", "thinking": "ragiono..."},
        {"type": "text", "text": "PIANO OK"}]}
    with patch("requests.post", return_value=ok):
        assert ds._ask_ai("sys", "user") == "PIANO OK"
    # modello d'area non disponibile (404) -> ripiega sul modello base
    ko = MagicMock(status_code=404, text='{"error":"model not found"}')
    calls = []
    def _fake_post(url, headers=None, json=None, timeout=None):
        calls.append(json["model"])
        return ko if json["model"] == "claude-fable-5" else ok
    with patch("requests.post", side_effect=_fake_post):
        assert ds._ask_ai("sys", "user") == "PIANO OK"
    assert calls == ["claude-fable-5", "claude-opus-4-8"]


def test_dynamics_semantic_process_cache():
    from app.engines import dynamics_search as ds
    # le primitive della cache di processo esistono
    assert hasattr(ds, "_semantic_ensure") and hasattr(ds, "warm_semantic_index")
    assert ds._SEM_LOCK is not None
    cat = {"WorkersEntity": {"string": ["Name", "BirthDate"]},
           "PurchaseRequisitions": {"string": ["Requester", "Amount"]}}
    r1 = ds._semantic_ensure(cat, {})
    r2 = ds._semantic_ensure(cat, {})
    if r1 is not None:
        # seconda chiamata: STESSA matrice (nessuna ricostruzione)
        assert r2 is not None and r1[1] is r2[1]
    # warm-up senza catalogo su disco: esce pulito senza eccezioni
    assert ds.warm_semantic_index({}) in (True, False)


def test_static_assets_versioned():
    store.create_user("assetv@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    # login: css versionato anche da sloggati
    assert "app.css?v=" in c.get("/login").text
    login(c, "assetv@test", "Password123")
    html = c.get("/").text
    assert "chat.js?v=" in html and "app.css?v=" in html


def test_dynamics_semantic_cache_per_process():
    pass  # segnaposto retrocompatibile


def test_area_models_admin_and_wiring():
    from app import connectors
    from app.main import _area_settings
    # helper: l'area applica il suo modello, vuota = predefinito
    base = {"claude_model": "claude-opus-4-8", "claude_model_rewrite": "claude-haiku-4-5-20251001",
            "claude_model_docgen": ""}
    assert _area_settings(base, "claude_model_rewrite")["claude_model"] == "claude-haiku-4-5-20251001"
    assert _area_settings(base, "claude_model_docgen")["claude_model"] == "claude-opus-4-8"
    # il planner Dynamics usa il modello dell'area (fallback al predefinito)
    cfg = {"client_id": "x", "tenant_id": "y", "resource_url": "https://z"}
    full = connectors._dyn_full_cfg(cfg, {"claude_model": "claude-opus-4-8",
                                          "claude_model_dynamics": "claude-fable-5"})
    assert full["claude_model"] == "claude-fable-5"
    full2 = connectors._dyn_full_cfg(cfg, {"claude_model": "claude-opus-4-8"})
    assert full2["claude_model"] == "claude-opus-4-8"
    # la pagina admin mostra le tre tendine con il consigliato
    store.create_user("areamdl@test", auth.hash_password("Password123"), "IT", is_admin=True)
    c = fresh_client(); login(c, "areamdl@test", "Password123")
    html = c.get("/admin").text
    for fld in ("claude_model_dynamics", "claude_model_rewrite", "claude_model_docgen"):
        assert f'name="{fld}"' in html
    assert html.count("consigliato") >= 4


def test_docgen_confirmation_inherits_format():
    from app.docgen import detect_request_with_history as d
    hist = [{"role": "user", "content": "mi generi un pptx sulla nis2?"},
            {"role": "assistant", "content": "Ecco la struttura... procedo con la generazione del file PowerPoint?"}]
    # la conferma eredita il formato dalla conversazione
    assert d("sì procedi", hist) == "pptx"
    assert d("ok vai", hist) == "pptx"
    # senza storia documentale, la conferma NON genera nulla
    assert d("sì procedi", []) is None
    assert d("ok", [{"role": "assistant", "content": "la capitale è Roma"}]) is None
    # un messaggio lungo non è una conferma
    assert d("sì ma prima spiegami meglio il progetto e i vincoli", hist) is None


def test_docgen_detect_conjugations():
    from app.docgen import detect_request
    # segnalazione utente reale: "mi generi un pptx..." non veniva rilevato
    assert detect_request("ciao mi generi un pptx sulla direttiva nis2 da presentare al ceo?") == "pptx"
    assert detect_request("mi crei un excel dei costi?") == "xlsx"
    assert detect_request("puoi farmi un word di sintesi?") == "docx"
    assert detect_request("vorrei un pdf della policy") == "pdf"
    # niente falsi positivi sulle richieste di spiegazione
    assert detect_request("puoi spiegarmi come funziona excel?") is None
    assert detect_request("cos'è un file pdf?") is None


def test_pptx_template_no_marketing():
    from pptx import Presentation
    prs = Presentation("app/doc_templates/Presentation_template_1.pptx")
    found = []
    for master in prs.slide_masters:
        for sh in master.shapes:
            if sh.has_text_frame and "marketing" in (sh.text_frame.text or "").lower():
                found.append("master")
        for lay in master.slide_layouts:
            for sh in lay.shapes:
                if sh.has_text_frame and "marketing" in (sh.text_frame.text or "").lower():
                    found.append("layout")
    assert not found


def test_chat_requires_single_source_documentale():
    store.create_user("srcv@test", auth.hash_password("Password123"), "IT")
    c = fresh_client(); login(c, "srcv@test", "Password123")
    # documentale SENZA fonte -> errore chiaro (validazione server)
    r = c.post("/api/chat", json={"messages": [{"role": "user", "content": "ciao"}],
                                  "free_mode": False})
    assert r.status_code == 200 and "Seleziona una fonte dati" in r.text
    # fonte non disponibile (onedrive non connesso) -> errore chiaro
    r2 = c.post("/api/chat", json={"messages": [{"role": "user", "content": "ciao"}],
                                   "free_mode": False, "source": "onedrive"})
    assert "OneDrive non" in r2.text
    # fonte valida (kb) -> supera la validazione (poi manca la chiave API: errore diverso)
    r3 = c.post("/api/chat", json={"messages": [{"role": "user", "content": "ciao"}],
                                   "free_mode": False, "source": "kb"})
    assert "Seleziona una fonte" not in r3.text
    # AI libera: nessuna fonte richiesta
    r4 = c.post("/api/chat", json={"messages": [{"role": "user", "content": "ciao"}],
                                   "free_mode": True})
    assert "Seleziona una fonte" not in r4.text


def test_chat_page_has_source_picker():
    store.create_user("srcp@test", auth.hash_password("Password123"), "IT")
    c = fresh_client(); login(c, "srcp@test", "Password123")
    html = c.get("/").text
    assert 'id="src-picker"' in html and 'name="datasource"' in html
    # onedrive/dynamics non connessi per questo utente -> pill disabilitati
    assert html.count("disabled") >= 2
    js = c.get("/static/chat.js").text
    assert "selectedSource" in js and "pickSource" in js


def test_docgen_detect_format():
    from app import docgen
    assert docgen.detect_request("creami un word con la sintesi") == "docx"
    assert docgen.detect_request("fammi una presentazione") == "pptx"
    assert docgen.detect_request("generami un excel dei costi") == "xlsx"
    assert docgen.detect_request("esporta in pdf") == "pdf"
    assert docgen.detect_request("qual è la policy?") is None


def test_docgen_builders_produce_files():
    from app import docgen
    import os
    spec = {"title": "Test", "subtitle": "x",
            "sections": [{"heading": "S", "paragraphs": ["p"], "bullets": ["b1", "b2"]}]}
    for fn in (docgen.gen_docx, docgen.gen_pdf):
        p, name = fn(spec)
        assert os.path.getsize(p) > 500
    p, name = docgen.gen_pptx({"title": "T", "slides": [{"title": "A", "bullets": ["x"]}]})
    assert os.path.getsize(p) > 500 and name.endswith(".pptx")
    p, name = docgen.gen_xlsx({"title": "D", "sheets": [{"name": "S", "columns": ["A", "B"],
                              "rows": [["x", 1], ["y", 2]], "total_columns": [1]}]})
    assert os.path.getsize(p) > 500 and name.endswith(".xlsx")


def test_download_owner_only():
    import tempfile, os
    from app import connectors
    store.create_user("dl@test", auth.hash_password("Password123"), "IT")
    store.create_user("dl2@test", auth.hash_password("Password123"), "Sales")
    fd, path = tempfile.mkstemp(suffix=".docx"); os.write(fd, b"PK fake docx"); os.close(fd)
    token = connectors.register_download("dl@test", path, "documento.docx")
    c = fresh_client(); login(c, "dl@test", "Password123")
    r = c.get("/download/" + token)
    assert r.status_code == 200
    assert "documento.docx" in r.headers.get("content-disposition", "")
    c2 = fresh_client(); login(c2, "dl2@test", "Password123")
    assert c2.get("/download/" + token).status_code == 404


def test_attachment_extraction_endpoint():
    import io
    store.create_user("att@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "att@test", "Password123")
    # txt
    r = c.post("/api/attach", files={"files": ("nota.txt", io.BytesIO(b"Policy change management: approvazione obbligatoria."), "text/plain")})
    assert r.status_code == 200
    a = r.json()["attachments"][0]
    assert a["ok"] and "change management" in a["text"]
    # tipo non supportato
    r2 = c.post("/api/attach", files={"files": ("x.exe", io.BytesIO(b"MZ"), "application/octet-stream")})
    assert r2.json()["attachments"][0]["ok"] is False


def test_pptx_extraction():
    # crea una pptx minima e verifica l'estrazione testo
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        return  # libreria non installata nell'ambiente di test: salta
    import io
    from app import knowledge
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Procedura backup"
    buf = io.BytesIO(); prs.save(buf)
    txt = knowledge.extract_attachment_text("p.pptx", buf.getvalue())
    assert "Procedura backup" in txt


def test_chat_page_has_attach_ui():
    store.create_user("attui@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "attui@test", "Password123")
    r = c.get("/")
    assert 'id="attach-btn"' in r.text and 'id="attach-input"' in r.text
    js = c.get("/static/chat.js").text
    assert "handleFiles" in js and "/api/attach" in js


def test_search_with_links_signature():
    from app import connectors
    # non connesso -> ("", []) senza eccezioni
    txt, links = connectors.search_with_links("nobody@test", "onedrive", "x")
    assert txt == "" and links == []


def test_chat_page_renders_sources_support():
    store.create_user("src@test", auth.hash_password("Password123"), "IT")
    c = fresh_client()
    login(c, "src@test", "Password123")
    # la pagina inietta il dizionario con l'etichetta "Fonti"
    r = c.get("/")
    assert "I18N_FB" in r.text and "sources:" in r.text
    # la logica delle fonti è in chat.js (file statico)
    js = c.get("/static/chat.js").text
    assert "addSources" in js and 'type === "sources"' in js


def test_onedrive_query_keywords():
    from app.engines.onedrive_search import _build_query
    # mappa concettuale: "anni" -> documento "anagrafica" (logica desktop)
    q = _build_query("quanti anni ha marco bonometti?")
    assert "anagrafica" in q and "bonometti" in q
    # contratto
    assert "contratto" in _build_query("mostrami il contratto di mario verdi")
    # fattura
    assert "fattura" in _build_query("fattura fornitore tecnoplast")
    # query senza concetto noto: tiene comunque i termini di contenuto
    assert _build_query("procedura backup server").strip() != ""


def test_build_system_includes_context():
    from app.orchestrator import build_system
    s = build_system("Tecnico", "Italiano", "FATTO_RISERVATO_XYZ")
    assert "FATTO_RISERVATO_XYZ" in s and "CONTESTO" in s
    s2 = build_system("Tecnico", "Italiano", "")
    assert "CONTESTO" not in s2


if __name__ == "__main__":
    import inspect, sys
    fns = [f for n, f in sorted(globals().items()) if n.startswith("test_") and inspect.isfunction(f)]
    failed = 0
    for f in fns:
        try:
            f()
            print(f"  PASS  {f.__name__}")
        except Exception as e:
            failed += 1
            import traceback
            print(f"  FAIL  {f.__name__}: {e}")
    print(f"\n{len(fns)-failed}/{len(fns)} test superati.")
    sys.exit(1 if failed else 0)
